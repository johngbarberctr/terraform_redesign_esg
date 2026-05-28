"""Generate ACI_Redesign_Executive_Summary.pptx.

Leadership-targeted, 12-slide companion to ACI_Redesign_Strategy.pptx.
Distinct from the strategy deck: assumes the audience does not (and does
not need to) know what an EPG, ESG, or AAEP is. Covers business case,
current vs target state, Design A (UCS-FI direct attach), the cutover
sequence, risks/rollback, current status, and decisions still needed.

Run:
    cd aci-redesign
    python3 build_executive_summary_pptx.py
Output:
    aci-redesign/ACI_Redesign_Executive_Summary.pptx
"""
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# --- Theme (matches ACI_Redesign_Strategy.pptx) --------------------------
NAVY = RGBColor(0x0B, 0x2E, 0x5C)
CISCO_BLUE = RGBColor(0x00, 0x5A, 0x9C)
ACCENT = RGBColor(0xF2, 0x99, 0x00)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
GREY = RGBColor(0x55, 0x5C, 0x66)
DARK = RGBColor(0x1C, 0x1F, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1D, 0x8A, 0x3E)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xE8, 0xA2, 0x00)
PANEL_LINE = RGBColor(0xDD, 0xE2, 0xEA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]
slides = []


# --- Helpers -------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=14, color=DARK,
                font="Calibri", line_spacing=1.2):
    """bullets accepts three forms per item:
       * "text"                     -> top-level bullet
       * ("text", int_indent)       -> bullet at that indent level
       * ("Heading", "body text")   -> bold heading bullet, body sub-bullet
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)

    def _emit(p, text, indent, *, bold=False):
        p.level = indent
        p.line_spacing = line_spacing
        bullet_char = "•" if indent == 0 else "◦"
        run = p.add_run()
        run.text = f"{bullet_char}  {text}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    para_idx = 0
    for item in bullets:
        if isinstance(item, tuple) and len(item) == 2:
            a, b = item
            if isinstance(b, int):
                p = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
                _emit(p, a, b)
                para_idx += 1
            else:
                # (heading, body) -> two paragraphs, body indented one level
                p = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
                _emit(p, a, 0, bold=True)
                para_idx += 1
                p = tf.add_paragraph()
                _emit(p, b, 1)
                para_idx += 1
        else:
            p = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
            _emit(p, item, 0)
            para_idx += 1
    return tb


def page_header(slide, title, subtitle=None, *, badge="EXECUTIVE SUMMARY"):
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.9), NAVY)
    add_rect(slide, Emu(0), Inches(0.9), SLIDE_W, Inches(0.05), ACCENT)
    add_text(slide, Inches(0.4), Inches(0.08), Inches(11.5), Inches(0.6),
             title, size=26, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.55), Inches(11.5), Inches(0.35),
                 subtitle, size=12, color=RGBColor(0xCF, 0xDA, 0xEA),
                 anchor=MSO_ANCHOR.MIDDLE)
    if badge:
        bw, bh = Inches(2.0), Inches(0.38)
        bx = SLIDE_W - bw - Inches(0.3)
        by = Inches(0.26)
        add_rect(slide, bx, by, bw, bh, ACCENT)
        add_text(slide, bx, by, bw, bh, badge, size=10, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def page_footer(slide, page_num, total,
                label="ACI Redesign — Executive Summary"):
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.35), SLIDE_W,
             Inches(0.35), NAVY)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.32),
             Inches(8.0), Inches(0.3), label,
             size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(2.5), SLIDE_H - Inches(0.32),
             Inches(2.1), Inches(0.3), f"{page_num} / {total}",
             size=10, color=WHITE, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)


def panel(slide, x, y, w, h, title, body_lines, *, accent=CISCO_BLUE):
    add_rect(slide, x, y, w, h, LIGHT, line=PANEL_LINE)
    add_rect(slide, x, y, Inches(0.18), h, accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.1),
             w - Inches(0.4), Inches(0.45), title,
             size=15, bold=True, color=NAVY)
    add_bullets(slide, x + Inches(0.3), y + Inches(0.55),
                w - Inches(0.4), h - Inches(0.65),
                body_lines, size=12)


def stat_card(slide, x, y, w, h, value, label, color=CISCO_BLUE):
    add_rect(slide, x, y, w, h, WHITE, line=PANEL_LINE)
    add_rect(slide, x, y, w, Inches(0.18), color)
    add_text(slide, x, y + Inches(0.4), w, Inches(0.95), value,
             size=44, bold=True, color=color, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + h - Inches(0.6), w, Inches(0.5), label,
             size=12, color=DARK, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def add_table(slide, x, y, w, h, headers, rows, *,
              col_widths=None, header_bg=NAVY, header_fg=WHITE,
              body_size=12):
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, c in enumerate(col_widths):
            table.columns[i].width = int(int(w) * c / total)
    # header
    for i, t in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Emu(36000)
        tf.margin_right = Emu(36000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = t
        run.font.name = "Calibri"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = header_fg
    # body
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                LIGHT if r_idx % 2 == 0 else WHITE)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Emu(36000)
            tf.margin_right = Emu(36000)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = value
            run.font.name = "Calibri"
            run.font.size = Pt(body_size)
            run.font.color.rgb = DARK
    return table_shape


# --- Slides --------------------------------------------------------------
def title_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, NAVY)
    # Big accent bar
    add_rect(s, Inches(0.0), Inches(2.7), SLIDE_W, Inches(0.18), ACCENT)
    # Title
    add_text(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.0),
             "ACI Redesign", size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(2.95), Inches(11.5), Inches(0.7),
             "Executive Summary",
             size=28, color=RGBColor(0xCF, 0xDA, 0xEA))
    # Sub-line
    add_text(s, Inches(0.8), Inches(3.95), Inches(11.5), Inches(0.5),
             "AEDCG + AEDCK   ·   2-VRF design   ·   "
             "Design A (UCS-FI direct attach)   ·   IaC end-to-end",
             size=16, color=WHITE)
    # Card row at the bottom
    cards = [
        ("39", "EPGs (consolidated from 266)"),
        ("2", "VRFs (down from 11)"),
        ("39", "BDs (down from 215)"),
        ("100%", "Terraform-managed"),
    ]
    cw = Inches(2.7)
    gap = Inches(0.3)
    cards_total = len(cards) * cw + (len(cards) - 1) * gap
    cx = (SLIDE_W - cards_total) / 2
    cy = Inches(5.4)
    for value, label in cards:
        add_rect(s, cx, cy, cw, Inches(1.4), WHITE, line=PANEL_LINE)
        add_rect(s, cx, cy, cw, Inches(0.15), ACCENT)
        add_text(s, cx, cy + Inches(0.25), cw, Inches(0.65), value,
                 size=34, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.1), cy + Inches(0.95),
                 cw - Inches(0.2), Inches(0.4), label,
                 size=11, color=DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw + gap
    # Footer line
    add_text(s, Inches(0.8), SLIDE_H - Inches(0.6),
             Inches(11.5), Inches(0.4),
             f"Prepared {date.today():%B %Y}   ·   "
             "Network Engineering   ·   Confidential, internal use",
             size=11, color=RGBColor(0xCF, 0xDA, 0xEA))
    slides.append(s)


def headline_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "What this is",
                "One-page view of why we are doing this and what changes")
    # Three columns: WHY / WHAT / IMPACT
    cols = [
        ("Why",
         CISCO_BLUE,
         [
            "Today's design is 11 VRFs, 215 bridge domains, 266 endpoint "
            "groups grown organically over years.",
            "Every change requires a long maintenance window because "
            "blast radius is unclear.",
            "Front-end Nexus 5K switches add a hop, a failure domain, "
            "and a license/support cost we can retire.",
            "There is no single source of truth — drift between APIC "
            "and Excel/runbooks is constant.",
         ]),
        ("What",
         NAVY,
         [
            "Collapse to 2 VRFs (internal + DMZ) with vzAny "
            "permit-all baseline; tighten with ESGs over time.",
            "Consolidate to 39 BDs / 39 EPGs grouped by function, "
            "not by VLAN number.",
            "Move ESXi traffic onto an APIC-managed VDS (per-fabric "
            "APCG-VDS1, APCK-VDS1).",
            "Connect UCS Fabric Interconnects directly to ACI leaves "
            "(Design A) — N5Ks decommissioned.",
            "All of it is described as YAML and applied via "
            "Terraform + GitLab CI.",
         ]),
        ("Impact",
         ACCENT,
         [
            "No IP renumbering. No VLAN renumbering for VM "
            "workloads. No application changes.",
            "Cutover is staged: policy push first (zero traffic "
            "impact), then physical re-cable in a maintenance window.",
            "Lab is already cut over and validated. Production is "
            "ready pending a change window and a few decisions.",
            "After cutover: faster, safer changes; smaller blast "
            "radius; automated drift detection.",
         ]),
    ]
    col_w = Inches(4.1)
    gap = Inches(0.15)
    x = Inches(0.5)
    y = Inches(1.2)
    h = Inches(5.6)
    for title, color, items in cols:
        add_rect(s, x, y, col_w, h, LIGHT, line=PANEL_LINE)
        add_rect(s, x, y, col_w, Inches(0.55), color)
        add_text(s, x, y, col_w, Inches(0.55), title,
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.2), y + Inches(0.7),
                    col_w - Inches(0.4), h - Inches(0.8),
                    items, size=12, line_spacing=1.25)
        x += col_w + gap
    # Closing band
    add_rect(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             NAVY)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             "Same applications. Same IPs. Cleaner network. "
             "Safer change path.",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    slides.append(s)


def current_vs_target_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Today vs. After",
                "The numeric and operational delta in one view")
    rows = [
        ["VRFs", "11", "2 (VRF-EUR + VRF-DMZ)",
         "Easier to reason about traffic flow"],
        ["Bridge Domains", "215", "39",
         "Functional grouping replaces ad-hoc per-VLAN BDs"],
        ["Endpoint Groups", "266", "39",
         "Consolidated by workload type"],
        ["Endpoint Security Groups", "0", "Future (vzAny first)",
         "Phased tightening of policy without re-cabling"],
        ["L3Outs", "13", "Unchanged (legacy)",
         "Out of scope for this cutover; "
         "managed by the existing ndo-terraform repo"],
        ["VMware integration",
         "VMM1 (single, manual VLAN tagging)",
         "Per-fabric APCG-VDS1 / APCK-VDS1 (dynamic)",
         "ACI manages VLANs, port-groups, host uplinks"],
        ["Front-end switching",
         "Nexus 5K vPC fronting UCS",
         "UCS FIs directly attached to ACI (Design A)",
         "One fewer device per path; N5Ks decommissioned"],
        ["Configuration source of truth",
         "GUI + spreadsheets",
         "YAML + Terraform + Git",
         "Reviewable, auditable, reproducible"],
        ["Change pipeline",
         "Manual, ticket-driven",
         "GitLab CI: validate → plan → manual deploy",
         "Plan reviewed before any APIC/NDO write"],
    ]
    add_table(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.4),
              ["Dimension", "Today", "After redesign", "Why it matters"],
              rows, col_widths=[1.2, 1.7, 1.9, 2.3], body_size=12)
    # Bottom footnote
    add_text(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
             "Counts sourced from production NDO schema "
             "AEDCE / AppProf-NetCentric and "
             "docs/reports/bd_mapping_analysis.txt.",
             size=11, color=GREY, align=PP_ALIGN.CENTER)
    slides.append(s)


def architecture_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Target Architecture (Design A)",
                "UCS Fabric Interconnects directly attached to ACI leaves")
    # Two stacked sites side-by-side, each showing FI-A/B → leaves → ESXi → VDS
    site_w = Inches(6.05)
    site_h = Inches(5.0)
    site_y = Inches(1.15)
    sites = [
        ("AEDCG", Inches(0.5),
         "Leaves 152 + 153   ·   APCG-VDS1   ·   PC_FI_A → eth1/6   ·   PC_FI_B → eth1/7"),
        ("AEDCK", Inches(6.78),
         "Leaves 119 + 191   ·   APCK-VDS1   ·   PC_FI_A → eth1/6   ·   PC_FI_B → eth1/7"),
    ]
    for name, x, sub in sites:
        add_rect(s, x, site_y, site_w, site_h, LIGHT, line=PANEL_LINE)
        add_rect(s, x, site_y, site_w, Inches(0.5), CISCO_BLUE)
        add_text(s, x, site_y, site_w, Inches(0.5),
                 f"Site {name}", size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), site_y + Inches(0.55),
                 site_w - Inches(0.3), Inches(0.4), sub,
                 size=10, color=GREY, align=PP_ALIGN.CENTER)
        # Stack: ACI leaves → port-channels → FIs → UCS chassis → ESXi → VDS
        layers = [
            ("ACI leaves (spine-attached, ESG/EPG enforcement)",
             NAVY, WHITE),
            ("Port-channels  PC_FI_A   ·   PC_FI_B   (mac-pin)",
             CISCO_BLUE, WHITE),
            ("UCS Fabric Interconnects   FI-A   ·   FI-B",
             ACCENT, WHITE),
            ("UCS B-series blade chassis", GREEN, WHITE),
            ("ESXi hosts  (VTEP, vmkernel mgmt on EPG-VHOST-MGMT)",
             RGBColor(0x55, 0x70, 0xA0), WHITE),
            ("APIC-managed VDS  (port-groups = redesign EPGs)",
             RGBColor(0x6B, 0x46, 0x9E), WHITE),
        ]
        ly = site_y + Inches(1.0)
        layer_h = Inches(0.55)
        gap = Inches(0.05)
        for text, fill, fg in layers:
            add_rect(s, x + Inches(0.25), ly,
                     site_w - Inches(0.5), layer_h, fill)
            add_text(s, x + Inches(0.25), ly,
                     site_w - Inches(0.5), layer_h, text,
                     size=11, bold=True, color=fg,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            ly += layer_h + gap
    # Bottom band: what NDO does
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3),
             Inches(0.7), NAVY)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3),
             Inches(0.35),
             "Tenant tree (2 VRFs, 39 BDs, 39 EPGs, contracts) "
             "is pushed once from NDO to BOTH sites simultaneously.",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3),
             Inches(0.35),
             "Each APIC owns its own per-fabric access policies "
             "and adopts its own pre-existing per-fabric VDS.",
             size=11, color=RGBColor(0xCF, 0xDA, 0xEA),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    slides.append(s)


def what_does_not_change_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "What does NOT change",
                "The redesign is non-disruptive on these dimensions")
    items_left = [
        ("IP addressing",
         "Every existing subnet keeps its IP. "
         "Bridge-domain consolidation absorbs subnets, never re-numbers them."),
        ("Application traffic",
         "vzAny permit-all in each VRF preserves all current "
         "intra-VRF flows. No application owner has to do anything."),
        ("VM workloads",
         "VMs land on the new APIC-managed VDS via a vMotion-equivalent "
         "port-group migration. No guest reconfiguration."),
        ("Tenant boundaries",
         "Tenant EUR is unchanged. The redesign sits inside it."),
    ]
    items_right = [
        ("North-south routing (L3Outs)",
         "The 13 existing L3Outs stay where they are, managed by the "
         "legacy ndo-terraform repo. Same external peers, same routes."),
        ("vCenter",
         "Same vCenter instance, same datacenter, same hosts. "
         "ACI adopts the existing per-fabric VDS rather than creating a new one."),
        ("Hostnames, DNS, AD",
         "All identity-side dependencies are untouched."),
        ("Operational ownership",
         "Same network, virtualisation, and UCS teams. Tooling changes "
         "internally; the operating model stays the same."),
    ]
    col_w = Inches(6.0)
    y = Inches(1.2)
    h = Inches(5.6)
    panel(s, Inches(0.5), y, col_w, h,
          "User-visible state is preserved", items_left, accent=GREEN)
    panel(s, Inches(6.85), y, col_w, h,
          "Adjacent systems are untouched", items_right, accent=CISCO_BLUE)
    slides.append(s)


def cutover_plan_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Cutover Plan",
                "Five gated stages — first three are reversible without traffic impact")
    rows = [
        ["1", "APIC access / fabric policies",
         "Push new VLAN pools, AAEPs, port-channels, leaf profiles, "
         "per-fabric VMM domains.",
         "None — additive only", "GREEN"],
        ["2", "Tenant tree via NDO",
         "Push 2 VRFs, 39 BDs, 39 EPGs, vzAny contracts. "
         "Bind EPGs to per-fabric VDSes.",
         "None — VMs stay on legacy port-groups until Stage 4",
         "GREEN"],
        ["3", "Static port bindings",
         "PATCH non-VM EPGs (load balancer, radio gateways, ESXi "
         "vmkernel) into NDO via curated JSON.",
         "None — bindings activate on the existing physical paths",
         "GREEN"],
        ["4", "UCS / vCenter physical move",
         "Re-cable FI uplinks from N5K to ACI leaves. "
         "Migrate VM port-groups to the new VDS.",
         "Brief blip per host pair as VDS uplinks move", "AMBER"],
        ["5", "Decommission",
         "After 24-hour soak: retire legacy VMM bindings and "
         "N5K front-end.",
         "None", "GREEN"],
    ]
    headers = ["Stage", "Action", "Detail", "Traffic impact", ""]
    add_table(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(4.6),
              headers, [r[:-1] + [""] for r in rows],
              col_widths=[0.4, 1.7, 4.0, 2.0, 0.4], body_size=11)
    # Status pills overlaid on column 5
    pill_x = Inches(12.43)
    pill_w = Inches(0.46)
    pill_h = Inches(0.32)
    py = Inches(1.55)
    row_h = Inches(0.85)
    pill_color = {"GREEN": GREEN, "AMBER": AMBER, "RED": RED}
    for r in rows:
        add_rect(s, pill_x, py + Inches(0.04),
                 pill_w, pill_h, pill_color[r[-1]])
        add_text(s, pill_x, py + Inches(0.04),
                 pill_w, pill_h, r[-1],
                 size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        py += row_h
    # Bottom: snapshots / gating
    add_rect(s, Inches(0.4), Inches(6.0), Inches(12.5),
             Inches(1.05), LIGHT, line=PANEL_LINE)
    add_rect(s, Inches(0.4), Inches(6.0), Inches(0.18),
             Inches(1.05), NAVY)
    add_text(s, Inches(0.7), Inches(6.05), Inches(12.0),
             Inches(0.4),
             "Pre-flight (T-7 days) and rollback are codified in the "
             "Production cutover runbook (aci-redesign/README.md)",
             size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12.0),
             Inches(0.65),
             "APIC + NDO + vCenter snapshots before Stage 1   ·   "
             "deploy jobs gated when:manual on main   ·   "
             "every stage has an explicit verify and rollback step",
             size=12, color=DARK)
    slides.append(s)


def risks_rollback_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Risks & Rollback",
                "Top three risks with named mitigation and rollback target")
    rows = [
        ["VLAN-collision on static bindings",
         "Multiple legacy EPGs consolidating into one redesign EPG can "
         "produce two VLANs on the same path/EPG, which APIC rejects.",
         "Pick a VLAN winner per (path, EPG) pair before the cutover. "
         "Run the bindings dump tool, review, edit, re-run.",
         "Stage 3 NDO undeploy + remove bindings — "
         "non-VMM EPGs revert to no-binding."],
        ["UCS / FI cabling mismatch",
         "Production FI uplinks are assumed on eth1/6 and eth1/7. "
         "If actual cabling differs, APIC won't program the bundle "
         "until YAML is corrected.",
         "Walk the cabling worksheet with UCS team in pre-flight. "
         "Per-leaf interface profiles are easy to amend in YAML.",
         "Stage 1 destroy / partial-fabric apply — additive only, "
         "rolls back in <60s."],
        ["NDO push collides with same-named APIC-local objects",
         "The first NDO deploy can detect existing tenant objects with "
         "matching names and need to absorb them.",
         "Plan reviewed before deploy. Deploy is manual (when:manual). "
         "Lab cutover already validated this path.",
         "Stage 2 NDO undeploy + terraform destroy on ndo/ root. "
         "Legacy IPv6 schema is in a separate state, untouched."],
    ]
    add_table(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(4.7),
              ["Risk", "Why it can happen", "Mitigation", "Rollback"],
              rows, col_widths=[1.5, 2.6, 3.0, 2.5], body_size=11)
    # Bottom band
    add_rect(s, Inches(0.4), Inches(6.05), Inches(12.5),
             Inches(1.0), LIGHT, line=PANEL_LINE)
    add_rect(s, Inches(0.4), Inches(6.05), Inches(0.18),
             Inches(1.0), GREEN)
    add_text(s, Inches(0.7), Inches(6.1), Inches(12.0),
             Inches(0.4), "Stages 1-3 are reversible with no traffic impact",
             size=14, bold=True, color=GREEN)
    add_text(s, Inches(0.7), Inches(6.45), Inches(12.0),
             Inches(0.55),
             "Only Stage 4 (physical re-cable + VDS uplink move) "
             "carries any traffic risk, and that risk is per-host blips, "
             "not site-wide outage. Worst case is restoring from APIC + "
             "NDO snapshots taken in pre-flight.",
             size=12, color=DARK)
    slides.append(s)


def status_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Status",
                "Where we are right now")
    # Three large status cards
    cards = [
        ("LAB", GREEN, "DONE",
         [
            "Per-fabric VMM domains (APCG-VDS1, APCK-VDS1) deployed "
            "and adopted from existing VDSes in shared vCenter.",
            "Tenant tree pushed via NDO: 2 VRFs, 39 BDs, 39 EPGs, "
            "vzAny contracts; deploy validated.",
            "39 port-groups visible in vCenter on each VDS, across "
            "3 ESXi hosts per site.",
            "All access/fabric policies match the netascode/nac-aci "
            "module's expected schema; no plan diffs.",
         ]),
        ("PRODUCTION", AMBER, "STAGED",
         [
            "Production Terraform root (apic-vmware-prod/) created "
            "and lints clean; mirrors lab structure.",
            "Production access policies committed for both fabrics: "
            "PC_FI_A/B, fi-static-vlan-pool (213 VLANs), "
            "leaf 152/153 + 119/191 split.",
            "GitLab CI jobs added: validate-aci-prod, plan-aci-prod, "
            "deploy-aci-prod (manual gate). Same for ndo-redesign root.",
            "Cutover runbook documented in aci-redesign/README.md.",
         ]),
        ("OUTSTANDING", RED, "DECISIONS NEEDED",
         [
            "Generate the production static-bindings JSON "
            "(read-only dump from prod NDO).",
            "Choose a VLAN-winner policy where multiple legacy EPGs "
            "consolidate onto one redesign EPG.",
            "Confirm UCS FI uplink ports (assumed eth1/6, eth1/7).",
            "Schedule the maintenance window for Stage 4 (physical "
            "FI re-cable + VDS uplink migration).",
         ]),
    ]
    cw = Inches(4.1)
    gap = Inches(0.15)
    x = Inches(0.5)
    y = Inches(1.2)
    h = Inches(5.6)
    for label, color, status, items in cards:
        add_rect(s, x, y, cw, h, LIGHT, line=PANEL_LINE)
        add_rect(s, x, y, cw, Inches(0.7), color)
        add_text(s, x, y, cw, Inches(0.4), label,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + Inches(0.35), cw, Inches(0.35), status,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.2), y + Inches(0.85),
                    cw - Inches(0.4), h - Inches(0.95),
                    items, size=12, line_spacing=1.25)
        x += cw + gap
    slides.append(s)


def benefits_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Operational Benefits After Cutover",
                "What changes about how we run the network day-to-day")
    rows = [
        ["Add a new EPG / workload",
         "Edit YAML → MR → CI plan → manual deploy",
         "Days → minutes; reviewable plan; no GUI clicks"],
        ["Onboard a new VM",
         "Land VM on the right port-group on the VDS",
         "Zero APIC change; ACI tracks endpoints automatically"],
        ["Tighten security",
         "Add ESGs and contracts incrementally; "
         "vzAny serves as the safety net",
         "Phased path; no flag-day"],
        ["Detect drift",
         "GitLab CI plan job runs on every change",
         "Discrepancies surface in MR review, not in production"],
        ["Recover from misconfig",
         "Revert MR → CI re-plans → manual deploy",
         "Rollback is a Git operation"],
        ["Disaster recovery",
         "Snapshots + Terraform state + Git history",
         "Same fabric configuration can be re-pushed to a "
         "rebuilt APIC in hours, not days"],
        ["Audit / compliance",
         "Every change is a Git commit + GitLab pipeline + "
         "Terraform state row",
         "Single chain of evidence per change"],
        ["Front-end switching cost",
         "N5Ks decommissioned at end of cutover",
         "License / support / power savings; one less failure "
         "domain in the data path"],
    ]
    add_table(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.6),
              ["Activity", "How it works after cutover", "Why it matters"],
              rows, col_widths=[1.5, 3.5, 3.0], body_size=12)
    slides.append(s)


def decisions_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Decisions Needed",
                "Inputs the cutover plan is waiting on")
    rows = [
        ["VLAN-winner policy",
         "Some redesign EPGs consolidate multiple legacy EPGs that "
         "were on different VLANs. APIC accepts only one VLAN per "
         "(path, EPG) pair.",
         "Network Engineering",
         "Three reasonable defaults: lowest VLAN wins, highest "
         "endpoint-count wins, or hand-curated. Recommend "
         "endpoint-count for least disruption."],
        ["Production maintenance window",
         "Stage 4 (physical FI re-cable, VDS uplink migration) "
         "needs a coordinated window across Network, UCS, "
         "and Virtualisation.",
         "Change Advisory Board",
         "Recommended: a single 4-hour window per site, two weeks "
         "apart so rollback signal is clear before site B."],
        ["UCS FI port mapping confirmation",
         "Production data files assume FI-A → eth1/6 and "
         "FI-B → eth1/7. UCS team to confirm or correct.",
         "UCS Engineering",
         "Quick check; YAML edit if different. No code change "
         "required."],
        ["L3Out scope",
         "Decision already made: L3Outs stay in the legacy "
         "ndo-terraform schema for this cutover. Revisit after "
         "stabilisation.",
         "Network Architecture",
         "Confirmed; documented in schema header and runbook."],
        ["ESG / micro-segmentation roadmap",
         "Phase 1 vzAny+permit-all is live. Phase 2 lift-and-shift ESGs "
         "(ESG-All-Internal-V2, ESG-All-DMZ-V2 in AppProf-AppCentric-V2) "
         "are in flight, applied APIC-direct via netascode/nac-aci@0.7.0 "
         "because the nac-ndo/mso provider does not model "
         "endpoint_security_groups. Phase 3 (per-zone ESGs with vCenter-tag "
         "selectors) is the follow-up design ticket.",
         "Security + Network",
         "Phase 2 keeps reachability identical (vzAny still permits all) "
         "and is fully reversible. Approve continued cadence."],
    ]
    add_table(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.7),
              ["Decision", "Why it's needed", "Owner", "Recommendation"],
              rows, col_widths=[1.6, 3.0, 1.3, 3.1], body_size=11)
    slides.append(s)


def timeline_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Indicative Timeline",
                "Right-to-left: each box is days, not weeks")
    # Horizontal timeline
    bar_y = Inches(3.4)
    bar_h = Inches(0.6)
    add_rect(s, Inches(0.5), bar_y, Inches(12.3), bar_h,
             RGBColor(0xE6, 0xEC, 0xF2))
    milestones = [
        ("Lab cutover", "DONE", GREEN, 0.04),
        ("Bindings JSON\n+ collision review", "T-7 days", AMBER, 0.20),
        ("Pre-flight\nsnapshots + plans", "T-3 days", AMBER, 0.36),
        ("Stage 1-3\npolicy + tenant + bindings", "T-0 morning",
         CISCO_BLUE, 0.55),
        ("Stage 4\nFI re-cable + VDS uplinks", "T-0 window",
         ACCENT, 0.74),
        ("Stage 5\ndecommission", "T+1 to T+7", NAVY, 0.92),
    ]
    box_w = Inches(1.85)
    box_h = Inches(1.5)
    for label, when, color, frac in milestones:
        cx = Inches(0.5) + Inches(12.3 * frac) - box_w / 2
        # Marker on bar
        add_rect(s, cx + box_w / 2 - Inches(0.1),
                 bar_y - Inches(0.1), Inches(0.2),
                 bar_h + Inches(0.2), color)
        # Box above bar
        by = bar_y - box_h - Inches(0.4)
        add_rect(s, cx, by, box_w, box_h, WHITE, line=PANEL_LINE)
        add_rect(s, cx, by, box_w, Inches(0.35), color)
        add_text(s, cx, by, box_w, Inches(0.35), when,
                 size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.1), by + Inches(0.4),
                 box_w - Inches(0.2), box_h - Inches(0.5),
                 label, size=11, color=DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Caption
    add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
             "Total active engineering effort during the change "
             "window: ~4-6 hours per site.",
             size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
             "Timeline assumes a single change-window per site, with "
             "AEDCG and AEDCK cut over two weeks apart so issues at "
             "site A surface before site B is touched.",
             size=12, color=GREY,
             align=PP_ALIGN.CENTER)
    # Soft-band at bottom
    add_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
             LIGHT, line=PANEL_LINE)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
             "Concrete dates set once VLAN-winner policy + UCS port "
             "confirmation are in.",
             size=12, color=GREY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    slides.append(s)


def ask_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Ask",
                "What we need from this audience")
    items = [
        ("Approve the cutover approach",
         "Single-window-per-site, Design A (UCS-FI direct attach), "
         "deploy gated when:manual."),
        ("Sponsor the maintenance window",
         "Two 4-hour windows, two weeks apart (AEDCG then AEDCK). "
         "Network, UCS, and Virtualisation owners on the bridge."),
        ("Sign off on VLAN-winner policy",
         "Recommend: highest endpoint-count wins for each "
         "consolidated EPG."),
        ("Confirm UCS FI port assumption",
         "FI-A → eth1/6, FI-B → eth1/7 on each leaf pair. "
         "UCS team to validate or adjust."),
        ("Acknowledge deferred items",
         "L3Out reorganisation is a follow-up, not part of this cutover. "
         "ESG-based micro-segmentation is in flight: Phase 2 lift-and-shift "
         "ESGs ship with this cutover (reachability-neutral); Phase 3 "
         "per-zone ESGs are the next design ticket."),
    ]
    panel(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
          "Decisions / approvals", items, accent=ACCENT)
    add_rect(s, Inches(0.5), Inches(6.85), Inches(12.3),
             Inches(0.4), NAVY)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3),
             Inches(0.4),
             "Everything else is built, lab-validated, "
             "in source control, and CI-gated.",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    slides.append(s)


def references_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Where to dig deeper",
                "Pointers for engineers and reviewers")
    rows = [
        ["Strategy deck (technical)",
         "aci-redesign/ACI_Redesign_Strategy.pptx",
         "27 slides covering the full design, BD mapping, "
         "phased security, deployment design decisions."],
        ["Production cutover runbook",
         "aci-redesign/README.md → 'Production cutover runbook' section",
         "Pre-flight, 5-stage sequence with verify steps, "
         "per-stage rollback table."],
        ["NDO schema (single source of truth)",
         "aci-redesign/data/nac-ndo/schema-aedce-v2.nac.yaml",
         "2 VRFs, 39 BDs, 39 EPGs, vzAny + 2 contracts (all -V2). "
         "Header documents all in-scope and out-of-scope items."],
        ["Production access policies",
         "aci-redesign/data/nac-aci-{aedcg,aedck}-prod/access-policies.nac.yaml",
         "Design A: PC_FI_A/B, fi-static-vlan-pool, leaf splits."],
        ["Terraform roots",
         "aci-redesign/{apic-vmware,ndo,apic-vmware-prod}/",
         "Lab APIC root, NDO redesign root, production APIC root. "
         "Each has its own Makefile, scripts, and CI jobs."],
        ["Bindings tooling",
         "aci-redesign/scripts/{dump_bindings.py,deploy_bindings.py}",
         "Read-only dump from prod NDO; PATCH static bindings into "
         "the redesign schema."],
        ["GitLab CI",
         ".gitlab-ci.yml",
         "validate-aci, plan-aci, deploy-aci (lab) + "
         "validate-aci-prod, plan-aci-prod, deploy-aci-prod (prod, "
         "manual gate) + validate/plan/deploy-ndo-redesign."],
    ]
    add_table(s, Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.5),
              ["Topic", "Path", "What you'll find"], rows,
              col_widths=[1.6, 3.4, 3.5], body_size=11)
    add_rect(s, Inches(0.4), Inches(6.85), Inches(12.5),
             Inches(0.4), NAVY)
    add_text(s, Inches(0.4), Inches(6.85), Inches(12.5),
             Inches(0.4),
             "Repository: terraform_redesign_esg "
             "(GitLab: root/terraform_redesign_esg)",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    slides.append(s)


# --- Build deck ---------------------------------------------------------
title_slide()
headline_slide()
current_vs_target_slide()
architecture_slide()
what_does_not_change_slide()
cutover_plan_slide()
risks_rollback_slide()
status_slide()
benefits_slide()
decisions_slide()
timeline_slide()
ask_slide()
references_slide()

total = len(slides)
for i, slide in enumerate(slides):
    if i == 0:
        continue
    page_footer(slide, i + 1, total)

out = ("/Users/johbarbe/DC/ACI/terraform-esg/aci-redesign/"
       "ACI_Redesign_Executive_Summary.pptx")
prs.save(out)
print(f"Saved {out} with {total} slides.")
