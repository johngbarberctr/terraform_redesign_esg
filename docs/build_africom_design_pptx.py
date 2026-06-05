"""Build AFRICOM_ACI_Design_Review.pptx

AFRICOM NIPR ACI design review deck — tenant AFR-DEL.Services.
Source data: CX Optimization doc (Sept 2025), CX Resiliency doc (March 2026),
and the ESG/NDO redesign strategy developed for this engagement.
No RCC-E / EUR tenant content included.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
NAVY    = RGBColor(0x0B, 0x2E, 0x5C)
BLUE    = RGBColor(0x00, 0x5A, 0x9C)
ACCENT  = RGBColor(0xF2, 0x99, 0x00)
LIGHT   = RGBColor(0xF4, 0xF6, 0xFA)
GREY    = RGBColor(0x55, 0x5C, 0x66)
DARK    = RGBColor(0x1C, 0x1F, 0x24)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x1D, 0x8A, 0x3E)
RED     = RGBColor(0xC0, 0x39, 0x2B)
AMBER   = RGBColor(0xE8, 0xA2, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

slides = []

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def rect(slide, x, y, w, h, fill, line=None, lw=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.shadow.inherit = False
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    return s


def txt(slide, x, y, w, h, text, *, sz=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font; r.font.size = Pt(sz)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, *, sz=13, color=DARK, ls=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    for i, item in enumerate(items):
        text, indent = (item[0], item[1]) if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = indent; p.line_spacing = ls
        r = p.add_run()
        r.text = ("•  " if indent == 0 else "◦  ") + text
        r.font.name = "Calibri"; r.font.size = Pt(sz)
        r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None, badge=None):
    rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.9), NAVY)
    rect(slide, Emu(0), Inches(0.9), SLIDE_W, Inches(0.05), ACCENT)
    txt(slide, Inches(0.4), Inches(0.08), Inches(11.5), Inches(0.6),
        title, sz=26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        txt(slide, Inches(0.4), Inches(0.55), Inches(11.5), Inches(0.35),
            subtitle, sz=12, color=RGBColor(0xCF, 0xDA, 0xEA),
            anchor=MSO_ANCHOR.MIDDLE)
    if badge:
        bw, bh = Inches(1.8), Inches(0.38)
        bx = SLIDE_W - bw - Inches(0.3)
        by = Inches(0.26)
        rect(slide, bx, by, bw, bh, ACCENT)
        txt(slide, bx, by, bw, bh, badge, sz=11, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, n, total):
    rect(slide, Emu(0), SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), NAVY)
    txt(slide, Inches(0.4), SLIDE_H - Inches(0.32), Inches(8),
        Inches(0.32), "AFRICOM NIPR — ACI Design Review", sz=10,
        color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.32),
        Inches(1.2), Inches(0.32), f"{n} / {total}",
        sz=10, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def table(slide, x, y, w, h, headers, rows, *,
          hbg=NAVY, hfg=WHITE, zebra=True, hsz=12, bsz=11, cw=None):
    nr, nc = len(rows) + 1, len(headers)
    shape = slide.shapes.add_table(nr, nc, x, y, w, h)
    tbl = shape.table
    if cw:
        total = sum(cw)
        for i, c in enumerate(cw):
            tbl.columns[i].width = int(w * c / total)
    for j, h_txt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = hbg
        cell.text_frame.clear()
        cell.text_frame.margin_left = cell.text_frame.margin_right = Emu(36000)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = h_txt
        r.font.bold = True; r.font.size = Pt(hsz)
        r.font.color.rgb = hfg; r.font.name = "Calibri"
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if (zebra and i % 2 == 0) else WHITE
            cell.text_frame.clear()
            cell.text_frame.margin_left = cell.text_frame.margin_right = Emu(36000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(bsz); r.font.color.rgb = DARK; r.font.name = "Calibri"
    return tbl


def panel(slide, x, y, w, h, title, body, *, accent=BLUE):
    rect(slide, x, y, w, h, LIGHT, line=RGBColor(0xDD, 0xE2, 0xEA))
    rect(slide, x, y, w, Inches(0.45), accent)
    txt(slide, x + Inches(0.15), y, w - Inches(0.2), Inches(0.45),
        title, sz=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    bullets(slide, x + Inches(0.1), y + Inches(0.5),
            w - Inches(0.2), h - Inches(0.55), body, sz=11)


def stat_card(slide, x, y, w, h, number, label, color):
    rect(slide, x, y, w, h, LIGHT, line=RGBColor(0xDD, 0xE2, 0xEA))
    rect(slide, x, y, Inches(0.12), h, color)
    txt(slide, x + Inches(0.2), y + Inches(0.08), w - Inches(0.3),
        Inches(0.7), number, sz=26, bold=True, color=color)
    txt(slide, x + Inches(0.2), y + Inches(0.82), w - Inches(0.3),
        Inches(0.4), label, sz=11, color=GREY)


def arrow_right(slide, x, y, w, h, color=NAVY):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, NAVY)
    rect(s, Emu(0), Inches(5.0), SLIDE_W, Inches(2.5), BLUE)
    rect(s, Emu(0), Inches(4.95), SLIDE_W, Inches(0.07), ACCENT)
    txt(s, Inches(0.8), Inches(1.5), Inches(11.8), Inches(1.2),
        "AFRICOM NIPR — ACI Design Review",
        sz=48, bold=True, color=WHITE)
    txt(s, Inches(0.8), Inches(2.8), Inches(11.8), Inches(0.7),
        "Tenant AFR-DEL.Services  |  Kelley & Del Din  |  Application-Centric Transition",
        sz=20, color=RGBColor(0xE5, 0xEE, 0xF8))
    rect(s, Inches(0.8), Inches(3.7), Inches(4.2), Inches(0.5), ACCENT)
    txt(s, Inches(0.8), Inches(3.7), Inches(4.2), Inches(0.5),
        "ESG-Driven Segmentation Strategy",
        sz=14, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.8), Inches(5.4), Inches(11.8), Inches(0.4),
        "Sources: CX Optimization Assessment (Sept 2025) · CX Resiliency Assessment (March 2026)",
        sz=13, color=WHITE)
    txt(s, Inches(0.8), Inches(6.1), Inches(11.8), Inches(0.4),
        "ACI 6.0(9e)  ·  NDO 4.4.3  ·  ND 3.2.2(m)  ·  2-Site Multi-Site Fabric",
        sz=12, color=RGBColor(0xAA, 0xBC, 0xD6))
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 2 — Agenda
# ---------------------------------------------------------------------------
def slide_agenda():
    s = prs.slides.add_slide(BLANK)
    header(s, "Agenda")
    items = [
        ("1  Current State",         "Fabric, tenant, and schema as assessed by Cisco CX"),
        ("2  Key Findings",           "Priority findings from Optimization & Resiliency reviews"),
        ("3  Problems to Solve",      "Why the current design limits application-centric growth"),
        ("4  Target Architecture",    "Revised VRF, BD, EPG, and template model"),
        ("5  Template Restructuring", "4-template target from 7 auto-generated templates"),
        ("6  Phased ESG Approach",    "vzAny → lift-and-shift → zone ESGs → explicit contracts"),
        ("7  Terraform Scope",        "NDO-focused: what TF manages vs what is already in place"),
        ("8  Execution Order",        "Dependency-correct deploy sequence"),
        ("9  Risks & Blockers",       "What can break and what must be decided first"),
        ("10 Questions to Resolve",   "5 decisions that gate the next phase"),
        ("11 Next Steps",             "Immediate actions"),
    ]
    col_w = Inches(5.9)
    lx, rx = Inches(0.5), Inches(7.0)
    top, rh = Inches(1.2), Inches(0.52)
    half = (len(items) + 1) // 2
    for idx, (k, v) in enumerate(items):
        col = 0 if idx < half else 1
        row = idx if col == 0 else idx - half
        x = lx if col == 0 else rx
        y = top + rh * row
        rect(s, x, y, Inches(0.22), rh - Inches(0.08), ACCENT)
        txt(s, x + Inches(0.3), y, col_w - Inches(0.3),
            Inches(0.26), k, sz=13, bold=True, color=NAVY)
        txt(s, x + Inches(0.3), y + Inches(0.25),
            col_w - Inches(0.3), Inches(0.27), v, sz=11, color=GREY)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 3 — Current State
# ---------------------------------------------------------------------------
def slide_current_state():
    s = prs.slides.add_slide(BLANK)
    header(s, "Current State — Fabric & Tenant Overview",
           "Data from Cisco CX Assessments (Sept 2025 / March 2026)")

    rows = [
        ["Sites",            "2  —  Kelley (NADE02) and Del Din (NAIT03)"],
        ["ACI Version",      "6.0(9e) on all leaves, spines, and APICs"],
        ["ND / NDO",         "Nexus Dashboard 3.2.2(m)  /  NDO 4.4.3"],
        ["Tenant",           "AFR-DEL.Services  (production)  +  666 PALE  (sandbox)"],
        ["Schema",           "AFRICOM NIPR  —  7 templates (auto-restructured on NDO 4.x upgrade)"],
        ["VRFs",             "9 at each site  (fabric-wide scale: 9 VRFs total)"],
        ["Bridge Domains",   "~32 at Kelley  /  ~28 at Del Din  in AFR-DEL.Services"],
        ["EPGs",             "~56 at Kelley  /  ~49 at Del Din  (fabric-wide)"],
        ["ESGs",             "None"],
        ["Contracts",        "vzAny in provider + consumer mode  →  VRF-wide permit-all (unenforced)"],
        ["L3Outs",           "Per-site L3Out to DC Firewall (vPC; redundant at Kelley & Del Din)"],
        ["VMM Integration",  "VMM domain connected to vCenter  —  several VDS currently offline"],
        ["ISN",              "MPLS Blackcore Transport  —  BGP EVPN full mesh across both spines per site"],
        ["ND Cluster",       "3-node VM cluster at Kelley only  —  no standby, no Del Din nodes"],
    ]
    table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.4),
          ["Attribute", "Value"], rows, cw=[1.1, 4.0], bsz=12)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 4 — Hardware Inventory
# ---------------------------------------------------------------------------
def slide_hardware():
    s = prs.slides.add_slide(BLANK)
    header(s, "Hardware & Software Inventory",
           "Current device matrix from Cisco CX Resiliency Assessment (March 2026)")

    hw_rows = [
        ["Kelley",   "APIC-SERVER-M3", "3", "NADE02NMP90/91/92", "Oct 2028"],
        ["Kelley",   "N9K-C9332C Spine", "2", "NADE02SP201/202", "Jul 2029"],
        ["Kelley",   "N9K-C93180YC-FX Leaf", "2", "NADE02LF101/102", "Jul 2029"],
        ["Del Din",  "APIC-SERVER-M3", "3", "NAIT03NMP90/91/92", "Oct 2028"],
        ["Del Din",  "N9K-C9504 Spine", "2", "NAIT03SP201/202", "N/A (not announced)"],
        ["Del Din",  "N9K-C93180YC-FX Leaf", "4", "NAIT03LF101/102/BL103/104", "Jul 2029"],
    ]
    table(s, Inches(0.5), Inches(1.1), Inches(9.5), Inches(4.0),
          ["Site", "Model", "Count", "Node IDs", "End of Security Patches"],
          hw_rows, cw=[0.7, 1.4, 0.4, 1.8, 1.0], bsz=11)

    panel(s, Inches(10.1), Inches(1.1), Inches(2.9), Inches(4.0),
          "Lifecycle notes",
          [
              "APIC-SERVER-M3 EOS announced May 2023. Security patches end Oct 2028.",
              "N9K-C93180YC-FX EOS Aug 2023. Security patches end Jul 2029.",
              "N9K-C9504 (Del Din spines) — no EOS announced.",
              "No per-node scalability risks identified by CX at current EPG/BD/VRF counts.",
          ], accent=AMBER)

    txt(s, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.35),
        "All firmware uniform at 6.0(9e) across both sites — no mixed-version risk.",
        sz=13, bold=True, color=GREEN)

    sw_rows = [
        ["Nexus Dashboard", "3.2.2(m)", "Not announced"],
        ["NDO (Orchestrator)", "4.4.3",   "Not announced"],
    ]
    table(s, Inches(0.5), Inches(5.7), Inches(6.0), Inches(1.3),
          ["Software", "Version", "End of Life"], sw_rows, cw=[1.5, 0.8, 1.0], bsz=11)

    panel(s, Inches(6.8), Inches(5.7), Inches(6.2), Inches(1.3),
          "Upgrade recommendation (from CX)",
          ["Upgrade to Nexus Dashboard 4.1.1 for consolidated services, improved GUI, and topology snapshots."],
          accent=BLUE)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 5 — Key CX Findings (Priority)
# ---------------------------------------------------------------------------
def slide_cx_findings():
    s = prs.slides.add_slide(BLANK)
    header(s, "Key Findings — Cisco CX Priority Summary",
           "Compiled from Optimization (Sept 2025) and Resiliency (March 2026) assessments")

    rows = [
        ["HIGH",   "Create transition plan to Application-Centric mode",
         "vzAny + Permit-Any contract effectively disables all contract enforcement. EPG-to-ESG migration required before segmentation is possible.",
         "Resiliency"],
        ["HIGH",   "Cleanup Firewall Policies",
         "Kelley and Del Din must share an identical ACL policy. Any drift between sites breaks L3Out failover — a policy change applied to one site must be applied to both.",
         "Resiliency"],
        ["HIGH",   "Cleanup BGP routing on Firewall & Fusion routers",
         "Sites receive different BGP prefix counts (Kelley 857, Del Din 872) causing asymmetric routing risk for any new subnets deployed.",
         "Resiliency"],
        ["HIGH",   "Cleanup VMM Integration",
         "Multiple VDS offline; credentials expired; management IP removed from ACI-managed VDS causing loss of ACI visibility and automation failures.",
         "Resiliency"],
        ["MED",    "Implement NDO HA (standby node)",
         "3-node cluster tolerates only 1 failure. No nodes at Del Din. NDO offline = no template deploys or schema changes.",
         "Resiliency"],
        ["MED",    "Cleanup NIPR Schema Templates",
         "7 templates after NDO 4.x auto-restructure. Stretched BD + EPG + Non-L2 should be combined into single 'Stretched Services' template.",
         "Resiliency"],
        ["MED",    "Cleanup Bridge Domain configurations",
         "BDs configured with BUM or Unicast Routing where not required, increasing unnecessary bandwidth utilization.",
         "Both"],
        ["MED",    "Bridge Domain subnets",
         "BD-Primary-02 has multiple subnets + DHCP policy: only primary subnet can relay DHCP. Must be evaluated for any multi-subnet consolidated BD.",
         "Both"],
        ["LOW",    "APIC backups — remote location",
         "All three APICs (Kelley, Del Din) backup to single server NADE02NMV07. Single point of failure for backup recovery.",
         "Resiliency"],
    ]
    table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0),
          ["Priority", "Finding", "Detail", "Source"],
          rows, cw=[0.5, 1.5, 4.5, 0.7], bsz=10, hsz=11)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 6 — Problems to Solve (Why redesign)
# ---------------------------------------------------------------------------
def slide_problems():
    s = prs.slides.add_slide(BLANK)
    header(s, "Problems to Solve",
           "Why the current design limits application-centric growth")

    problems = [
        ("vzAny is unenforced",
         RED,
         [
             "vzAny in provider + consumer mode on the production VRF means ALL contracts are bypassed.",
             "A Permit-Any EPG-level contract was added as a 'transition' — this confirms enforcement is off.",
             "ESG contracts written today silently do nothing until vzAny is removed.",
             "This is the #1 design risk — must be resolved before ESG segmentation has any value.",
         ]),
        ("9 VRFs — reason unknown",
         AMBER,
         [
             "Legacy VRFs likely created per-mission, per-office, or per-compliance zone.",
             "The firewall already enforces real security boundaries at the L3Out.",
             "If VRFs don't carry overlapping IP space, consolidation may be possible.",
             "Must audit VRF names and L3Out topology before any consolidation recommendation.",
         ]),
        ("No ESGs — no path to micro-segmentation",
         BLUE,
         [
             "Without ESGs, every new security requirement needs a new EPG or a VRF.",
             "ESGs provide a policy overlay on top of the existing EPG/BD network model.",
             "They enable zero-trust segmentation without IP changes or endpoint disruption.",
             "Phase 2 lift-and-shift ESGs are observation-only — safe to deploy immediately.",
         ]),
        ("Schema templates out of best practice",
         GREY,
         [
             "7 templates after NDO 4.x auto-restructure — more than needed for 2 sites.",
             "VRF is in a shared stretched template, not its own template (NDO 4.x best practice).",
             "Stretched BD, Stretched EPG, and Stretched Non-L2 should be one 'Stretched Services' template.",
             "CX recommended 4-template target: VRF / Stretched Services / Kelley Unique / Del Din Unique.",
         ]),
    ]
    col_w = Inches(5.9)
    for i, (title, color, body) in enumerate(problems):
        row, col = divmod(i, 2)
        x = Inches(0.5) + col * (col_w + Inches(0.3))
        y = Inches(1.15) + row * Inches(2.7)
        panel(s, x, y, col_w, Inches(2.55), title, body, accent=color)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 7 — Target Architecture
# ---------------------------------------------------------------------------
def slide_target():
    s = prs.slides.add_slide(BLANK)
    header(s, "Target Architecture — Application-Centric Redesign",
           "Tenant AFR-DEL.Services  |  2 sites  |  ESG-driven segmentation")

    rows = [
        ["Tenant",         "AFR-DEL.Services (unchanged)"],
        ["Schema",         "New schema alongside existing AFRICOM NIPR  (e.g. AFR-SERVICES-V2)"],
        ["Templates",      "4:  VRF  /  Stretched-Services  /  Kelley-Unique  /  DelDin-Unique"],
        ["VRFs",           "Reduce from 9  →  2 (internal + DMZ)  — pending IP overlap audit"],
        ["Bridge Domains", "Consolidate into functional names matching security zones (~20–30 target)"],
        ["EPGs",           "1:1 with BDs  —  descriptive names replacing numeric/VLAN-based names"],
        ["Contracts",      "vzAny permit-all initially  →  tightened via ESG-to-ESG contracts (Phase 4)"],
        ["ESGs (Phase 2)", "One ESG per VRF grouping all EPGs  —  observation only, no policy change"],
        ["ESGs (Phase 3)", "Per-zone ESGs via vCenter tag selectors  (aci-zone, aci-tier)"],
        ["L3Outs",         "Keep dedicated per-VRF L3Out per site  (~4 total)  —  already best practice"],
        ["VMM",            "Stabilize existing VMM domain  —  service account, management IP isolation"],
        ["Naming",         "-V2 suffix on all new tenant objects for safe coexistence during parallel run"],
    ]
    table(s, Inches(0.5), Inches(1.1), Inches(7.8), Inches(5.8),
          ["Attribute", "Target"], rows, cw=[1.1, 3.2], bsz=12)

    # Right column: stat cards
    cards = [
        ("9 → 2",    "VRFs (pending audit)",      BLUE),
        ("~32 → ~20","Bridge Domains (estimated)", GREEN),
        ("0 → 2",    "ESGs (Phase 2)",             ACCENT),
        ("0",        "IP Address Changes",          NAVY),
    ]
    cx, cy, cw, ch = Inches(8.8), Inches(1.1), Inches(2.2), Inches(1.35)
    for i, (num, lbl, color) in enumerate(cards):
        row, col = divmod(i, 2)
        x = cx + col * (cw + Inches(0.2))
        y = cy + row * (ch + Inches(0.15))
        stat_card(s, x, y, cw, ch, num, lbl, color)

    panel(s, Inches(8.6), Inches(4.3), Inches(4.4), Inches(2.6),
          "What does NOT change",
          [
              "All existing IP addresses and subnets — zero readdressing.",
              "Firewall enforces DMZ/internal boundary — unchanged.",
              "Physical fabric, leaf profiles, AAEP — already configured.",
              "ACI/NDO version — already on target (6.0(9e), NDO 4.4.3).",
              "L3Out topology — dedicated per-VRF is already best practice.",
          ], accent=GREEN)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 8 — VRF Investigation
# ---------------------------------------------------------------------------
def slide_vrf_investigation():
    s = prs.slides.add_slide(BLANK)
    header(s, "VRF Investigation — Before Any Consolidation",
           "Must answer these questions before recommending a target VRF count")

    txt(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.4),
        "The current 9 VRFs need to be understood before consolidation. "
        "Pull the NDO schema export — VRF names reveal the intent.",
        sz=14, color=DARK)

    questions = [
        ("Why do these VRFs exist?",
         BLUE,
         [
             "Are they per-mission-system, per-office, or per-compliance-zone?",
             "Each VRF likely has its own L3Out context on the firewall — "
             "confirm with the firewall team before removing any VRF.",
             "DoD environments often have VRFs tied to classification levels "
             "or specific command boundaries — these may NOT be consolidatable.",
         ]),
        ("Is consolidation safe?",
         AMBER,
         [
             "Check for overlapping IP space across VRFs — any two VRFs sharing "
             "a prefix CANNOT be merged into a single VRF.",
             "Run: show ip route vrf <name> on each border leaf and compare.",
             "If any two VRFs have the same subnet, the merge would create a "
             "routing conflict that breaks both workloads.",
         ]),
        ("What is the firewall topology?",
         RED,
         [
             "Each VRF likely maps to a separate firewall context (security zone).",
             "Collapsing VRFs means collapsing firewall contexts — "
             "that is a firewall project, not an ACI project.",
             "Minimum viable path: keep as many VRFs as the firewall has contexts "
             "that must remain separated.",
         ]),
        ("What is the realistic target?",
         GREEN,
         [
             "Best case (all internal, no IP overlap): 2 VRFs — Internal + DMZ.",
             "Likely case (some mandatory separation): 3–4 VRFs.",
             "Conservative case (compliance-driven): keep existing count; "
             "use ESGs for micro-segmentation within each VRF.",
             "ESGs work in ANY VRF count — don't block Phase 2/3 on this decision.",
         ]),
    ]
    col_w = Inches(5.9)
    for i, (title, color, body) in enumerate(questions):
        row, col = divmod(i, 2)
        x = Inches(0.5) + col * (col_w + Inches(0.3))
        y = Inches(1.65) + row * Inches(2.6)
        panel(s, x, y, col_w, Inches(2.45), title, body, accent=color)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 9 — Template Restructuring
# ---------------------------------------------------------------------------
def slide_templates():
    s = prs.slides.add_slide(BLANK)
    header(s, "NDO Template Restructuring",
           "From 7 auto-generated templates to 4 best-practice templates")

    # Current state
    txt(s, Inches(0.5), Inches(1.1), Inches(5.9), Inches(0.38),
        "Current State (7 templates)", sz=15, bold=True, color=NAVY)
    current = [
        ("Stretched VRF",       "Contains VRF",                            AMBER),
        ("Stretched BD",        "Stretched Bridge Domains",                 AMBER),
        ("Stretched EPG",       "Stretched EPGs",                           AMBER),
        ("Stretched Non-L2",    "Non-L2-stretched BDs/EPGs",                AMBER),
        ("Kelley Unique",       "Site-local objects",                       BLUE),
        ("Del Din Unique",      "Site-local objects",                       BLUE),
    ]
    th = Inches(0.6)
    for i, (name, detail, color) in enumerate(current):
        y = Inches(1.55) + i * (th + Inches(0.06))
        rect(s, Inches(0.5), y, Inches(5.9), th, LIGHT,
             line=RGBColor(0xDD, 0xE2, 0xEA))
        rect(s, Inches(0.5), y, Inches(0.12), th, color)
        txt(s, Inches(0.75), y, Inches(2.5), th, name,
            sz=12, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(3.3), y, Inches(3.0), th, detail,
            sz=11, color=GREY, anchor=MSO_ANCHOR.MIDDLE)

    # Arrow
    arrow_right(s, Inches(6.55), Inches(3.5), Inches(0.5), Inches(0.45), ACCENT)

    # Target state
    txt(s, Inches(7.2), Inches(1.1), Inches(5.7), Inches(0.38),
        "Target State (4 templates — CX Recommended)", sz=15, bold=True, color=GREEN)
    target = [
        ("VRF",
         "VRFs, Filters, Contracts, Service Chains. "
         "VRF always gets its own template in NDO 4.x — "
         "prevents circular dependency on deploy/undeploy.",
         GREEN),
        ("Stretched Services",
         "All stretched BDs and EPGs (collapses current "
         "Stretched BD + Stretched EPG + Stretched Non-L2 into one). "
         "Deployed to both Kelley and Del Din.",
         BLUE),
        ("Kelley Unique",
         "Site-local BDs, EPGs, and L3Outs not used at Del Din. "
         "Assigned only to Kelley site.",
         NAVY),
        ("Del Din Unique",
         "Site-local BDs, EPGs, and L3Outs not used at Kelley. "
         "Assigned only to Del Din site.",
         NAVY),
    ]
    th2 = Inches(1.2)
    for i, (name, detail, color) in enumerate(target):
        y = Inches(1.55) + i * (th2 + Inches(0.1))
        rect(s, Inches(7.2), y, Inches(5.7), th2, LIGHT,
             line=RGBColor(0xDD, 0xE2, 0xEA))
        rect(s, Inches(7.2), y, Inches(0.12), th2, color)
        txt(s, Inches(7.45), y + Inches(0.08), Inches(5.3), Inches(0.4),
            name, sz=13, bold=True, color=DARK)
        txt(s, Inches(7.45), y + Inches(0.45), Inches(5.3), Inches(0.68),
            detail, sz=11, color=GREY)

    rect(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.45), NAVY)
    txt(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.45),
        "Deploy order: VRF template first → Stretched Services → Kelley Unique → Del Din Unique. "
        "Undeploy in reverse. NDO 4.x enforces this — circular deps cause deploy failures.",
        sz=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 10 — Phased ESG Approach
# ---------------------------------------------------------------------------
def slide_esg_phases():
    s = prs.slides.add_slide(BLANK)
    header(s, "Phased ESG Approach",
           "Start permissive, classify, then tighten — no disruption at any phase")

    phases = [
        ("Phase 1", "vzAny permit-all\n(already deployed)",
         "VRF-wide permit-all. All EPGs communicate freely. "
         "Safety net that makes every subsequent phase reversible.",
         "Complete", GREEN),
        ("Phase 2", "Lift-and-shift ESGs\n(deploy now)",
         "One ESG per VRF groups all EPGs. vzAny stays. "
         "ESGs are observation-only — classify endpoints, change nothing. "
         "Verify endpoint counts match EPG sums.",
         "Ready to deploy", BLUE),
        ("Phase 3", "Per-zone ESGs\n(requires tag scheme owner)",
         "Split ESGs by vCenter tag selectors (aci-zone, aci-tier). "
         "Map legacy VRF zone names to ESG names. "
         "vzAny still in place — no traffic risk.",
         "Blocked on governance", AMBER),
        ("Phase 4", "Explicit contracts\n(requires ADM data)",
         "Replace vzAny with ESG-to-ESG contracts covering only "
         "documented flows. Remove vzAny from DMZ VRF first, "
         "then Internal VRF after 30-day zero-hit confirmation.",
         "Future", GREY),
    ]
    bw = Inches(2.95); gap = Inches(0.2)
    sx = Inches(0.5); y = Inches(1.3)
    for i, (name, title, desc, status, color) in enumerate(phases):
        x = sx + i * (bw + gap)
        rect(s, x, y, bw, Inches(3.8), LIGHT,
             line=RGBColor(0xDD, 0xE2, 0xEA))
        rect(s, x, y, bw, Inches(0.55), color)
        txt(s, x, y, bw, Inches(0.55), name, sz=14, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.2), y + Inches(0.65), bw - Inches(0.4),
            Inches(1.0), title, sz=13, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.2), y + Inches(1.6), bw - Inches(0.4),
            Inches(1.8), desc, sz=10, color=DARK, align=PP_ALIGN.CENTER)
        rect(s, x + Inches(0.5), y + Inches(3.38),
             bw - Inches(1.0), Inches(0.32), color)
        txt(s, x + Inches(0.5), y + Inches(3.38),
            bw - Inches(1.0), Inches(0.32), status, sz=10,
            bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            arrow_right(s, x + bw + Inches(0.02), y + Inches(1.75),
                        gap - Inches(0.04), Inches(0.4), ACCENT)

    # vzAny removal note
    rect(s, Inches(0.5), Inches(5.35), Inches(12.3), Inches(1.8),
         LIGHT, line=RGBColor(0xDD, 0xE2, 0xEA))
    rect(s, Inches(0.5), Inches(5.35), Inches(0.15), Inches(1.8), RED)
    txt(s, Inches(0.8), Inches(5.4), Inches(12.0), Inches(0.38),
        "Critical: vzAny removal prerequisites (Phase 4 only)",
        sz=14, bold=True, color=RED)
    bullets(s, Inches(0.8), Inches(5.8), Inches(12.0), Inches(1.25),
            [
                "All 'must-talk' flows from Application Dependency Mapping are covered by ESG-to-ESG contracts.",
                "Zero vzAny hits in flow logs over 30 days for any flow not already covered by an explicit contract.",
                "Maintenance window scheduled. Remove DMZ VRF vzAny first, Internal VRF after stabilization period.",
            ], sz=12)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 11 — Terraform Scope
# ---------------------------------------------------------------------------
def slide_terraform_scope():
    s = prs.slides.add_slide(BLANK)
    header(s, "Terraform Scope for AFRICOM",
           "APIC fabric is already configured — focus is the NDO tenant policy layer")

    # Left: NOT in scope
    panel(s, Inches(0.5), Inches(1.1), Inches(5.9), Inches(5.6),
          "Already configured — do NOT Terraform",
          [
              "APIC access/fabric policies (leaf profiles, AAEP, VPC policy groups).",
              "VLAN pools and physical/VMM domains (already deployed).",
              "MCP Instance Policy (already set with production key — do not overwrite).",
              "Leaf switch profiles and interface selectors.",
              "ISN / BGP EVPN spine peering.",
              "Firewall L3Out BGP peers.",
              "",
              "These are managed by the existing APIC directly or by a separate "
              "Terraform root that is not part of this engagement.",
          ], accent=GREY)

    # Right: IN scope
    panel(s, Inches(6.6), Inches(1.1), Inches(6.4), Inches(2.6),
          "Terraform: NDO root (ndo/)  — primary deliverable",
          [
              "New schema (e.g. AFR-SERVICES-V2) with 4-template structure.",
              "VRFs, contracts (vzAny), BDs with descriptive names, EPGs (1:1 with BDs).",
              "EPG-to-VMM-domain bindings per site.",
              "deploy_templates=false — operator clicks Deploy in NDO UI.",
              "All objects carry -V2 suffix for coexistence with AFRICOM NIPR schema.",
          ], accent=NAVY)

    panel(s, Inches(6.6), Inches(3.85), Inches(6.4), Inches(2.85),
          "APIC-direct (minimal, ESG layer only)",
          [
              "AppProf-AppCentric-V2  —  the ANP container for ESGs.",
              "ESG-All-Internal-V2  —  EPG selectors for all internal EPGs.",
              "ESG-All-DMZ-V2  —  EPG selectors for all DMZ EPGs.",
              "Applied AFTER NDO deploy (EPGs must exist on APIC first).",
              "Uses nac-aci@0.7.0 wrapper — nac-ndo/mso provider has no ESG resource.",
              "Future: consolidate into NDO root when nac-ndo adds ESG support.",
          ], accent=BLUE)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 12 — Execution Order
# ---------------------------------------------------------------------------
def slide_execution_order():
    s = prs.slides.add_slide(BLANK)
    header(s, "Terraform Execution Order",
           "Dependency-correct deploy sequence for AFRICOM")

    steps = [
        ("Pre-flight", GREY,
         "Before any Terraform runs",
         [
             "NDO Orchestrator app enabled on ND (token → POST /api/v1/licensetier).",
             "source .env  (NDO_URL, credentials).",
             "Verify NDO connectivity: make auth-check in ndo/.",
             "APIC snapshot on both Kelley and Del Din (rollback target).",
             "NDO backup: Operations → Backup.",
         ]),
        ("Step 1", BLUE,
         "Existing AFRICOM NIPR schema\n(already deployed — verify, don't touch)",
         [
             "Confirm AFRICOM NIPR schema is deployed at both sites in NDO.",
             "Note the filter 'Any' location in that schema — the V2 schema cross-references it.",
             "Alternatively: define 'Any' filter natively in the V2 schema to remove the cross-schema dependency.",
         ]),
        ("Step 2", NAVY,
         "NDO root: new V2 schema",
         [
             "terraform init + terraform plan + terraform apply in ndo/.",
             "Creates: schema AFR-SERVICES-V2, 4 templates, VRFs, contracts, BDs, EPGs.",
             "Nothing reaches APICs yet — deploy_templates=false.",
             "NDO UI: deploy VRF template first, then Stretched Services, then site-unique templates.",
         ]),
        ("Step 3", GREEN,
         "ESG layer (APIC-direct, re-apply after NDO deploy)",
         [
             "EPGs must exist on APIC before ESG selectors can resolve.",
             "Apply nac-aci wrapper with tenant-eur-esgs.nac.yaml for both sites.",
             "Verify: APIC GUI → Tenants → AFR-DEL.Services → AppProf-AppCentric-V2 → Endpoint counts.",
         ]),
        ("Step 4", ACCENT,
         "Static port bindings (Python scripts)",
         [
             "dump_bindings.py to read existing AFRICOM NIPR bindings.",
             "deploy_bindings.py to PATCH the V2 EPGs.",
             "NDO UI: re-deploy Stretched Services template to push bindings to fabric.",
         ]),
    ]
    step_h = Inches(1.1)
    x, w = Inches(0.5), Inches(12.3)
    for i, (step, color, title, body) in enumerate(steps):
        y = Inches(1.1) + i * (step_h + Inches(0.05))
        rect(s, x, y, Inches(0.8), step_h, color)
        txt(s, x, y, Inches(0.8), step_h, step, sz=13, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, x + Inches(0.85), y, w - Inches(0.85), step_h, WHITE,
             line=RGBColor(0xDD, 0xE2, 0xEA))
        txt(s, x + Inches(1.0), y + Inches(0.05), Inches(3.0), Inches(0.45),
            title, sz=12, bold=True, color=NAVY)
        # body as inline text
        body_text = "  ·  ".join(body)
        txt(s, x + Inches(1.0), y + Inches(0.5), w - Inches(1.2),
            Inches(0.55), body_text, sz=10, color=DARK)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 13 — Risks & Blockers
# ---------------------------------------------------------------------------
def slide_risks():
    s = prs.slides.add_slide(BLANK)
    header(s, "Risks & Blockers",
           "What can break and what must be decided before starting")

    rows = [
        ["vzAny supersedes all contracts",
         "Contracts written against ESGs silently do nothing while vzAny is active. Phase 4 has zero effect until vzAny is removed.",
         "Document the prerequisite checklist (E4). Do not remove vzAny until all flows are covered by explicit contracts.",
         "HIGH"],
        ["DHCP relay / multi-subnet BDs",
         "Consolidating many legacy BDs into one new BD puts multiple subnets on a single BD. Only the primary subnet can unicast for DHCP leases.",
         "Audit each consolidated BD: identify DHCP-served subnets. Split any BD where more than one subnet needs DHCP relay.",
         "HIGH"],
        ["VMM integration instability",
         "Multiple VDS offline at time of CX assessment. Credentials expired, management IP removed from ACI-managed VDS.",
         "Fix before deploying new ESGs: (1) vCenter management on standard vSwitch, (2) non-expiring service account, (3) change-control process between ACI and vCenter teams.",
         "HIGH"],
        ["Asymmetric BGP routing",
         "Sites receive different prefix counts (857 vs 872). New V2 BD subnets advertised into BGP may route asymmetrically, breaking firewall state.",
         "Verify all sites receive identical BGP prefix set. Confirm next-hop-self on firewall BGP peers before deploying any new subnets.",
         "MED"],
        ["VRF consolidation IP overlap",
         "If any two legacy VRFs share a subnet, they cannot be merged. Merging would create a routing conflict.",
         "Pull routing tables from each VRF before any consolidation plan. IP overlap = hard block on that VRF merge.",
         "MED"],
        ["Cross-schema filter dependency",
         "V2 schema references 'Any' filter from AFRICOM NIPR schema. If AFRICOM NIPR is restructured, V2 terraform plan fails.",
         "Define 'Any' filter natively in the V2 schema. Removes the dependency entirely.",
         "LOW"],
        ["nac-ndo has no ESG support",
         "ESG layer must be managed APIC-direct (not via NDO). This creates two-headed ownership — NDO doesn't see ESG membership.",
         "All ESG changes go through Terraform only, never via APIC GUI. Track nac-ndo upstream for ESG support.",
         "LOW"],
    ]
    table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0),
          ["Risk", "Impact", "Mitigation", "Level"],
          rows, cw=[1.4, 2.3, 2.6, 0.45], bsz=10)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 14 — Questions to Resolve
# ---------------------------------------------------------------------------
def slide_questions():
    s = prs.slides.add_slide(BLANK)
    header(s, "Questions to Resolve Before Starting",
           "5 decisions that gate the next phase — these cannot be answered by Terraform")

    questions = [
        ("1",
         "Why does AFRICOM have 9 VRFs?",
         BLUE,
         [
             "Pull the NDO schema export and read the VRF names — they will tell the story.",
             "Determine if any VRF boundaries are tied to firewall contexts, "
             "compliance zones, or mission-system ownership.",
             "Check for IP space overlap across VRFs — this is the hard constraint "
             "on consolidation, not the VRF count itself.",
             "Result: either a target VRF count, or a recommendation to keep the count "
             "and use ESGs for segmentation within each VRF.",
         ]),
        ("2",
         "Who owns the vCenter tag scheme?",
         AMBER,
         [
             "Phase 3 (per-zone ESGs via tag selectors) cannot start without a named owner "
             "for the aci-zone / aci-tier tag categories in vCenter.",
             "Network team defines the allowed values; vCenter team applies them to VMs.",
             "Must be agreed before any Phase 3 work begins — this is a people decision, not a tooling one.",
         ]),
        ("3",
         "Is NDI available for Application Dependency Mapping?",
         RED,
         [
             "Phase 4 contracts must cover every active flow before vzAny is removed.",
             "Without flow data, contract authoring is guesswork — and a missed flow causes an outage.",
             "NDI is the only available tool. Confirm licensing and whether flow analysis "
             "is enabled for the AFR-DEL.Services VRFs.",
         ]),
        ("4",
         "Does any compliance requirement prohibit VRF-wide permit-all?",
         NAVY,
         [
             "vzAny + permit-all is the migration scaffold — it is explicitly temporary.",
             "If any STIG, DISA, or audit requirement prohibits it, Phase 1 cannot deploy as designed.",
             "Alternative: contract every EPG pair before cutover, but this requires completing "
             "Phase 4 ADM work before Phase 1, inverting the dependency.",
         ]),
        ("5",
         "What is the long-term plan for AFRICOM NIPR schema?",
         GREEN,
         [
             "Option A: V2 schema eventually owns all objects; AFRICOM NIPR is retired (preferred).",
             "Option B: AFRICOM NIPR stays alive for L3Outs; V2 coexists indefinitely.",
             "Option B means the -V2 naming and two-schema ownership is permanent — harder to audit.",
             "This decision sets the decommission timeline and defines when the project is 'done'.",
         ]),
    ]
    qh = Inches(1.15)
    for i, (num, title, color, body) in enumerate(questions):
        y = Inches(1.1) + i * (qh + Inches(0.05))
        rect(s, Inches(0.5), y, Inches(0.6), qh, color)
        txt(s, Inches(0.5), y, Inches(0.6), qh, num, sz=20, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, Inches(1.15), y, Inches(11.65), qh, WHITE,
             line=RGBColor(0xDD, 0xE2, 0xEA))
        txt(s, Inches(1.3), y + Inches(0.06), Inches(3.0), Inches(0.36),
            title, sz=12, bold=True, color=NAVY)
        body_text = "  ·  ".join(body)
        txt(s, Inches(1.3), y + Inches(0.42), Inches(11.2),
            Inches(0.65), body_text, sz=10, color=DARK)
    slides.append(s)


# ---------------------------------------------------------------------------
# Slide 15 — Next Steps
# ---------------------------------------------------------------------------
def slide_next_steps():
    s = prs.slides.add_slide(BLANK)
    header(s, "Next Steps", "Immediate actions — in priority order")

    steps = [
        ("Immediate (this week)",
         GREEN,
         [
             "Pull NDO schema export for AFRICOM NIPR — read VRF names and L3Out topology.",
             "Verify all sites receive identical BGP prefix counts from firewall.",
             "Confirm VMM integration is stable: service account, management IP, VDS status.",
             "Identify whether any STIG/compliance rule prohibits VRF-wide permit-all.",
         ]),
        ("Short-term (2–4 weeks)",
         BLUE,
         [
             "Run IP overlap check across all 9 VRFs — determine consolidation feasibility.",
             "Name a vCenter tag scheme owner (joint: Network team + vCenter team).",
             "Draft 4-template NDO schema structure for review.",
             "Begin ndo/ Terraform root — write schema YAML for VRF + Stretched Services templates.",
             "Confirm NDI is licensed and enable flow analysis on AFR-DEL.Services VRFs.",
         ]),
        ("Medium-term (Phase 2 deploy)",
         NAVY,
         [
             "Apply ndo/ Terraform root — deploy V2 schema alongside AFRICOM NIPR.",
             "NDO UI: deploy templates in correct order (VRF → Stretched Services → site-unique).",
             "Apply APIC-direct ESG layer (AppProf-AppCentric-V2 + 2 ESGs).",
             "Verify endpoint counts in ESG operational view match EPG sum.",
             "Capture Phase 2 baseline: endpoint CSV export per ESG for Phase 3 comparison.",
         ]),
    ]
    col_w = Inches(3.9); gap = Inches(0.2)
    for i, (title, color, body) in enumerate(steps):
        x = Inches(0.5) + i * (col_w + gap)
        rect(s, x, Inches(1.1), col_w, Inches(5.5), LIGHT,
             line=RGBColor(0xDD, 0xE2, 0xEA))
        rect(s, x, Inches(1.1), col_w, Inches(0.5), color)
        txt(s, x, Inches(1.1), col_w, Inches(0.5), title, sz=13,
            bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        bullets(s, x + Inches(0.2), Inches(1.7),
                col_w - Inches(0.4), Inches(4.7), body, sz=12)

    rect(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4), NAVY)
    txt(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
        "Phase 2 ESG deploy is safe to start now — it is observation-only and fully reversible.",
        sz=13, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    slides.append(s)


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------
slide_title()
slide_agenda()
slide_current_state()
slide_hardware()
slide_cx_findings()
slide_problems()
slide_target()
slide_vrf_investigation()
slide_templates()
slide_esg_phases()
slide_terraform_scope()
slide_execution_order()
slide_risks()
slide_questions()
slide_next_steps()

total = len(slides)
for i, slide in enumerate(slides):
    if i == 0:
        continue
    footer(slide, i + 1, total)

out = "/Users/johbarbe/DC/ACI/sac-johbarbe-AFRICOM-terraform-esg-nac-ndo/docs/AFRICOM/AFRICOM_ACI_Design_Review.pptx"
prs.save(out)
print(f"Saved: {out}  ({total} slides)")
