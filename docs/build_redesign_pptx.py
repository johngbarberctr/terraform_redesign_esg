"""Generate ACI_Redesign_Strategy.pptx from aci-redesign documentation.

Builds a stand-alone deck covering:
 * why the redesign is needed
 * target 2-VRF architecture
 * BD consolidation & naming rules
 * lab (greenfield) and production (brownfield) migration phases
 * risks, mitigations, and decommission plan
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Theme ---------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x2E, 0x5C)       # primary header
CISCO_BLUE = RGBColor(0x00, 0x5A, 0x9C)  # secondary
ACCENT = RGBColor(0xF2, 0x99, 0x00)      # accent orange
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)       # panel bg
GREY = RGBColor(0x55, 0x5C, 0x66)        # subtitle
DARK = RGBColor(0x1C, 0x1F, 0x24)        # body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1D, 0x8A, 0x3E)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xE8, 0xA2, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# --- Helpers -------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=14, color=DARK,
                bullet_color=None, font="Calibri", line_spacing=1.15):
    """bullets: list of str or (str, int_indent) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, indent = item
        else:
            text, indent = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = indent
        p.line_spacing = line_spacing
        bullet_char = "•" if indent == 0 else "◦"
        run = p.add_run()
        run.text = f"{bullet_char}  {text}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if indent == 0:
            run.font.bold = False
    return tb


def page_header(slide, title, subtitle=None, *, badge=None):
    # Top banner
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.9), NAVY)
    add_rect(slide, Emu(0), Inches(0.9), SLIDE_W, Inches(0.05), ACCENT)
    add_text(slide, Inches(0.4), Inches(0.08), Inches(11.5), Inches(0.6),
             title, size=26, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.55), Inches(11.5), Inches(0.35),
                 subtitle, size=12, color=RGBColor(0xCF, 0xDA, 0xEA),
                 anchor=MSO_ANCHOR.MIDDLE)
    if badge:
        bw, bh = Inches(1.6), Inches(0.38)
        bx = SLIDE_W - bw - Inches(0.3)
        by = Inches(0.26)
        shp = add_rect(slide, bx, by, bw, bh, ACCENT)
        add_text(slide, bx, by, bw, bh, badge, size=11, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def page_footer(slide, page_num, total, label="ACI Redesign Strategy"):
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32),
             NAVY)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.32), Inches(8),
             Inches(0.32), label, size=10, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.32),
             Inches(1.2), Inches(0.32), f"{page_num} / {total}",
             size=10, color=WHITE, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)


def add_table(slide, x, y, w, h, headers, rows, *,
              header_bg=NAVY, header_fg=WHITE, zebra=True,
              header_size=12, body_size=11, col_widths=None):
    nrows = len(rows) + 1
    ncols = len(headers)
    shape = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    tbl = shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for j, h_text in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.text_frame.clear()
        cell.text_frame.margin_left = Emu(36000)
        cell.text_frame.margin_right = Emu(36000)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h_text
        r.font.bold = True
        r.font.size = Pt(header_size)
        r.font.color.rgb = header_fg
        r.font.name = "Calibri"
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (LIGHT if (zebra and i % 2 == 0)
                                        else WHITE)
            cell.text_frame.clear()
            cell.text_frame.margin_left = Emu(36000)
            cell.text_frame.margin_right = Emu(36000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(body_size)
            r.font.color.rgb = DARK
            r.font.name = "Calibri"
    return tbl


def section_label(slide, x, y, w, h, text, color=CISCO_BLUE):
    add_rect(slide, x, y, Inches(0.08), h, color)
    add_text(slide, x + Inches(0.15), y, w - Inches(0.15), h, text,
             size=14, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)


def panel(slide, x, y, w, h, title, body_lines, *, accent=CISCO_BLUE):
    add_rect(slide, x, y, w, h, LIGHT, line=RGBColor(0xDD, 0xE2, 0xEA))
    add_rect(slide, x, y, w, Inches(0.45), accent)
    add_text(slide, x + Inches(0.15), y, w - Inches(0.2), Inches(0.45),
             title, size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(slide, x + Inches(0.1), y + Inches(0.5),
                w - Inches(0.2), h - Inches(0.55), body_lines,
                size=11, color=DARK)


# --- Slide builders ------------------------------------------------------
slides = []


def title_slide():
    s = prs.slides.add_slide(BLANK)
    # Background gradient (two stacked rectangles for simplicity)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Emu(0), Inches(5.2), SLIDE_W, Inches(2.3), CISCO_BLUE)
    # Accent bar
    add_rect(s, Emu(0), Inches(5.15), SLIDE_W, Inches(0.06), ACCENT)
    # Title
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.8), Inches(1.2),
             "EUR Tenant ACI Redesign", size=54, bold=True,
             color=WHITE, font="Calibri")
    add_text(s, Inches(0.8), Inches(3.1), Inches(11.8), Inches(0.7),
             "From 11 VRFs / 266 EPGs to 2 VRFs / 39 EPGs with ESG-Ready Segmentation",
             size=22, color=RGBColor(0xE5, 0xEE, 0xF8))
    # Accent pill
    add_rect(s, Inches(0.8), Inches(4.0), Inches(3.4), Inches(0.5), ACCENT)
    add_text(s, Inches(0.8), Inches(4.0), Inches(3.4), Inches(0.5),
             "Strategy & Migration Process", size=15, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # Footer block
    add_text(s, Inches(0.8), Inches(5.6), Inches(11.8), Inches(0.4),
             "Consolidate routing · Adopt descriptive naming · Introduce ESGs · Preserve every IP",
             size=16, color=WHITE)
    add_text(s, Inches(0.8), Inches(6.3), Inches(11.8), Inches(0.4),
             "Project: terraform-esg / aci-redesign   |   Tenant: EUR   |   Schema: AFRICOM",
             size=12, color=RGBColor(0xCC, 0xD8, 0xEA))
    add_text(s, Inches(0.8), Inches(6.8), Inches(11.8), Inches(0.3),
             "Prepared from REDESIGN_OVERVIEW.md and aci-redesign/README.md",
             size=11, color=RGBColor(0xAA, 0xBC, 0xD6))
    slides.append(s)


def agenda_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Agenda", "What this deck covers and the order we walk through it")
    items = [
        ("1  Why redesign", "Problems with today's 11-VRF / 266-EPG layout"),
        ("2  Target architecture", "2 VRFs, 39 BDs/EPGs, ESGs, vzAny"),
        ("3  Key design decisions", "Why 2 VRFs, why L2 'DMZ' stays in VRF-AFR-DEL.Services-V2, no IP changes, why -V2"),
        ("4  BD consolidation", "How 215 legacy BDs map into 39 functional BDs"),
        ("5  Naming & object model", "Descriptive naming conventions"),
        ("6  Phased security model", "vzAny → single ESG → zone ESGs → tight contracts"),
        ("7  Lab migration (greenfield)", "4 phases - build from scratch"),
        ("8  Production migration (brownfield)", "7 phases - live coexistence"),
        ("9  VMM / VLAN strategy", "VMM pool 3501-3967 + per-fabric VMMs (APCG-VDS1, APCK-VDS1)"),
        ("10  Risks & mitigations", "Where traffic loss can happen and how we contain it"),
        ("11  What gets deployed", "Concrete object list from Terraform apply"),
        ("12  Decommission plan", "30 legacy BDs to remove"),
        ("13  Deployment architecture", "How YAML becomes ACI config via nac-aci + MCP sidecar"),
        ("14  Automation design decisions", "Four non-obvious choices that protect the pipeline"),
        ("15  Secrets strategy", "Env → CI → Vault progression, same TF_VAR_* names"),
        ("16  Local developer workflow", "What an operator actually types"),
        ("17  CI/CD pipeline", "GitLab stages: validate → plan → deploy → destroy"),
        ("18  Operations & troubleshooting", "auth-check, diagnose, common failure modes"),
        ("19  Next steps", "What to do after this deck"),
    ]
    # Two-column list
    col_w = Inches(5.9)
    left_x = Inches(0.5)
    right_x = Inches(7.0)
    top = Inches(1.2)
    row_h = Inches(0.48)
    half = (len(items) + 1) // 2
    for idx, (k, v) in enumerate(items):
        col = 0 if idx < half else 1
        row = idx if col == 0 else idx - half
        x = left_x if col == 0 else right_x
        y = top + row_h * row
        add_rect(s, x, y, Inches(0.22), row_h - Inches(0.08), ACCENT)
        add_text(s, x + Inches(0.3), y, col_w - Inches(0.3),
                 Inches(0.24), k, size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(0.22),
                 col_w - Inches(0.3), Inches(0.26), v,
                 size=11, color=GREY)
    slides.append(s)


def executive_summary_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Executive Summary",
                "One-page view of the redesign and why it matters")

    # Left narrative
    add_text(s, Inches(0.5), Inches(1.15), Inches(7.4), Inches(0.45),
             "The Goal", size=18, bold=True, color=NAVY)
    add_bullets(
        s, Inches(0.5), Inches(1.6), Inches(7.4), Inches(3.2),
        [
            "Collapse 11 legacy VRFs into 2 routing domains: VRF-AFR-DEL.Services-V2 (internal) and VRF-DMZ-V2 (proxy).",
            "Replace 215 numeric bridge domains and 266 EPGs with 39 descriptive, function-named BDs/EPGs that match the IPv6 RCC model.",
            "Introduce Endpoint Security Groups (ESGs) so segmentation can be enforced by policy instead of by VRF routing isolation.",
            "Keep every existing IPv4 subnet and gateway address in place - no host, VM, or switch is readdressed.",
            "Start permissive (vzAny permit-all) and tighten progressively using ESGs and targeted contracts.",
        ],
        size=13, color=DARK,
    )

    add_text(s, Inches(0.5), Inches(4.85), Inches(7.4), Inches(0.4),
             "Why now", size=18, bold=True, color=NAVY)
    add_bullets(
        s, Inches(0.5), Inches(5.25), Inches(7.4), Inches(1.9),
        [
            "11-VRF model added operational complexity without meaningful security (firewall already enforces real boundaries).",
            "Numeric naming (BD-V0372, EPG-V0572) hides function and slows incident response.",
            "No ESGs means no realistic path to micro-segmentation at scale.",
        ],
        size=13, color=DARK,
    )

    # Right stat cards
    cards = [
        ("11 → 2", "VRFs", CISCO_BLUE),
        ("215 → 39", "Bridge Domains", GREEN),
        ("266 → 39", "EPGs", ACCENT),
        ("13 → ~4", "L3Outs (prod)", RED),
        ("0", "IP Address Changes", NAVY),
        ("2", "New ESGs (Phase 2)", CISCO_BLUE),
    ]
    card_x = Inches(8.2)
    card_y = Inches(1.15)
    cw = Inches(2.3)
    ch = Inches(1.4)
    for i, (num, lbl, color) in enumerate(cards):
        row, col = divmod(i, 2)
        x = card_x + col * (cw + Inches(0.2))
        y = card_y + row * (ch + Inches(0.2))
        add_rect(s, x, y, cw, ch, LIGHT,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, x, y, Inches(0.12), ch, color)
        add_text(s, x + Inches(0.2), y + Inches(0.1), cw - Inches(0.3),
                 Inches(0.7), num, size=26, bold=True, color=color)
        add_text(s, x + Inches(0.2), y + Inches(0.85), cw - Inches(0.3),
                 Inches(0.4), lbl, size=12, color=GREY)
    slides.append(s)


def current_state_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Current State - What's Deployed Today",
                "EUR tenant / AFRICOM schema in production APIC")
    headers = ["Attribute", "Value"]
    rows = [
        ["Tenant", "EUR"],
        ["Schema", "AFRICOM"],
        ["VRFs (11)",
         "EUR-E, EUR-AIS, EUR-AIM, EUR-AIV, EUR-AIZ, EUR-AIG, EUR-AIP, "
         "EUR-AOV-UC-DMZ, EUR-ARMY-ENT-SVR-DMZ, EUR-GSN-Test, EUR-E catch-all"],
        ["Bridge Domains", "215  (numeric: BD-V0005 → BD-V2205)"],
        ["EPGs", "266  (numeric: EPG-V0005 → EPG-V2205)"],
        ["Contracts", "Per-VRF individual contracts (no vzAny)"],
        ["ESGs", "None"],
        ["L3Outs", "13"],
        ["VMM Domain", "VMM1"],
    ]
    add_table(s, Inches(0.5), Inches(1.15), Inches(7.8), Inches(4.0),
              headers, rows, col_widths=[1.1, 3.0])

    # Problems panel
    panel(
        s, Inches(8.6), Inches(1.15), Inches(4.3), Inches(5.8),
        "Problems With the Current Design",
        [
            "VRFs created for segmentation, but the firewall already enforces real security boundaries.",
            "Numeric BD/EPG names hide function - operators must cross-reference alias tables.",
            "266 EPGs are hard to audit, contract, and change.",
            "No ESGs = no realistic path to micro-segmentation.",
            "L3Out sprawl (13) complicates routing and changes.",
        ],
        accent=RED,
    )

    add_text(s, Inches(0.5), Inches(5.3), Inches(7.8), Inches(1.5),
             "Takeaway: the fabric works, but its object model fights operators every day. "
             "A flatter, name-driven model closes the gap without touching a single IP address.",
             size=13, bold=True, color=NAVY)
    slides.append(s)


def target_state_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Target State - 2-VRF V2 Redesign",
                "VRF-AFR-DEL.Services-V2 + VRF-DMZ-V2, 39 BDs, 39 EPGs, 2 ESGs, vzAny permit-all start")
    headers = ["Attribute", "Target"]
    rows = [
        ["Tenant", "EUR (unchanged)"],
        ["Schema / Template", "AFRICOM-V2 / Tenant_EUR_V2"],
        ["VRFs", "2  -  VRF-AFR-DEL.Services-V2 (internal) + VRF-DMZ-V2 (proxy segments)"],
        ["Bridge Domains",
         "39 descriptive  -  BD-AD-V2, BD-APP-SVR-V2, BD-CFG-MGMT-V2, ..."],
        ["EPGs", "39 (1:1 with BDs)  -  EPG-AD-V2, EPG-APP-SVR-V2, ..."],
        ["Contracts",
         "vzAny permit-all (Any_VRF-AFR-DEL.Services-V2, Any_VRF-DMZ-V2) - tightened via per-zone ESGs in Phase 3"],
        ["ESGs",
         "ESG-All-Internal-V2 (selects all 36 EPGs in AppProf-NetCentric-V2) + "
         "ESG-All-DMZ-V2 (selects all 3 EPGs in AppProf-DMZ-V2)  -  "
         "Phase 2 in flight; APIC-direct via nac-aci@0.7.0 (nac-ndo/mso provider lacks endpoint_security_groups)"],
        ["App Profiles",
         "3 ANPs: AppProf-NetCentric-V2 (NDO), AppProf-DMZ-V2 (NDO), AppProf-AppCentric-V2 (APIC-direct, holds the 2 ESGs)"],
        ["IPv6 (VRF-RCC)",
         "Unchanged today  -  folds into BD-*-V2 as additional subnets in a later wave"],
    ]
    add_table(s, Inches(0.5), Inches(1.15), Inches(7.8), Inches(4.2),
              headers, rows, col_widths=[1.1, 3.0])

    panel(
        s, Inches(8.6), Inches(1.15), Inches(4.3), Inches(4.2),
        "Why this shape (and why -V2)",
        [
            "VRF-AFR-DEL.Services-V2 consolidates all internal traffic into a single routing domain.",
            "VRF-DMZ-V2 keeps proxy segments (139.139.x.x) routing-isolated from internal.",
            "-V2 suffix lets AFRICOM-V2 coexist with the legacy AFRICOM schema in tenant EUR (ACI requires unique names per tenant).",
            "-V2 is generational, not address-family - same BDs will carry IPv4 + IPv6 once dual-stack wave lands.",
            "vzAny makes initial cut-over safe and reversible.",
            "ESGs give a classification layer today and a policy layer tomorrow.",
        ],
        accent=GREEN,
    )

    # BD counting band
    band_y = Inches(5.5)
    add_rect(s, Inches(0.5), band_y, Inches(12.3), Inches(1.5),
             LIGHT, line=RGBColor(0xDD, 0xE2, 0xEA))
    add_text(s, Inches(0.7), band_y + Inches(0.08), Inches(11.9),
             Inches(0.35), "Subnet consolidation across the 215 legacy BDs",
             size=13, bold=True, color=NAVY)
    tiles = [
        ("22", "BDs carrying real IPv4 subnets (110 subnets total)", CISCO_BLUE),
        ("17", "Placeholder BDs (IPv6-only categories)", ACCENT),
        ("14", "L2-only BDs (gateway on firewall)", GREEN),
        ("30", "Decommission (20 dead + 4 deprecated + 6 test)", RED),
    ]
    tile_w = Inches(2.9)
    for i, (n, lbl, c) in enumerate(tiles):
        x = Inches(0.7) + i * (tile_w + Inches(0.2))
        add_rect(s, x, band_y + Inches(0.5), tile_w, Inches(0.9),
                 WHITE, line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, x, band_y + Inches(0.5), Inches(0.08), Inches(0.9), c)
        add_text(s, x + Inches(0.15), band_y + Inches(0.55),
                 Inches(0.7), Inches(0.4), n, size=20, bold=True,
                 color=c)
        add_text(s, x + Inches(0.85), band_y + Inches(0.55),
                 tile_w - Inches(1.0), Inches(0.85), lbl,
                 size=10, color=GREY)
    slides.append(s)


def architecture_diagram_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Target Architecture Diagram",
                "Tenant EUR - two VRFs, vzAny per VRF, ESGs grouping all EPGs")

    # Tenant container
    tx, ty, tw, th = Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.8)
    add_rect(s, tx, ty, tw, th, WHITE,
             line=RGBColor(0xCF, 0xD6, 0xE0))
    add_rect(s, tx, ty, tw, Inches(0.4), NAVY)
    add_text(s, tx + Inches(0.15), ty, Inches(4), Inches(0.4),
             "Tenant: EUR", size=14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + Inches(4), ty, tw - Inches(4.2), Inches(0.4),
             "Filter: Any (cross-ref AFRICOM/VRF_Template/Any)  "
             "|  Contracts: Any_VRF-AFR-DEL.Services-V2, Any_VRF-DMZ-V2  "
             "|  Schema: AFRICOM-V2 / Tenant_EUR_V2",
             size=10, color=RGBColor(0xCF, 0xDA, 0xEA),
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    # VRF-EUR container
    ex, ey, ew, eh = tx + Inches(0.2), ty + Inches(0.6), Inches(7.7), Inches(5.0)
    add_rect(s, ex, ey, ew, eh, LIGHT,
             line=RGBColor(0xCF, 0xD6, 0xE0))
    add_rect(s, ex, ey, ew, Inches(0.45), CISCO_BLUE)
    add_text(s, ex + Inches(0.15), ey, ew, Inches(0.45),
             "VRF-AFR-DEL.Services-V2  -  Internal (IPv4 today, dual-stack target)   "
             "|   vzAny: Any_VRF-AFR-DEL.Services-V2 (permit-all)",
             size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    # BDs panel inside VRF-EUR
    bdx = ex + Inches(0.2)
    bdy = ey + Inches(0.55)
    bdw = ew - Inches(0.4)
    bdh = Inches(2.9)
    add_rect(s, bdx, bdy, bdw, bdh, WHITE,
             line=RGBColor(0xDD, 0xE2, 0xEA))
    add_text(s, bdx + Inches(0.1), bdy + Inches(0.05),
             bdw - Inches(0.2), Inches(0.3),
             "36 Bridge Domains  (22 with IPv4 subnets · 14 placeholders)",
             size=11, bold=True, color=NAVY)
    bd_list = [
        "BD-AD", "BD-APP-SVR (32)", "BD-CFG-MGMT (19)",
        "BD-WEB-SVR (11)", "BD-DB-SVR (8)", "BD-VVOIP-MGMT (9)",
        "BD-MECM (6)", "BD-BACKUP-SVR (4)", "BD-FILE-SVR (3)",
        "BD-SYSLOG (3)", "BD-LB (3)", "BD-ACAS-SCANNERS (2)",
        "BD-NMS", "BD-PATCH", "BD-PKI-SRV", "BD-VHOST-MGMT",
        "BD-VVOIP-PROXY", "BD-ADM-DCO",
        "BD-DNS-MGMT (L2)", "BD-GEF-MGMT (L2)",
        "BD-ACAS-MGMT*", "BD-ADFS*", "BD-C2C-SCANNERS*",
        "BD-DHCP-SVR*", "BD-E911-SVR*", "BD-FMWR-SVR*",
        "BD-LMR*", "BD-NAC*", "BD-OCSP*", "BD-PRINT-SVR*",
        "BD-RCC-DCO*", "BD-RCC-DNS*", "BD-RCC-SVR*",
        "BD-RCC-UNIX*", "BD-SMTP-SVR*", "BD-SYSMAN*",
    ]
    cols = 4
    chip_w = (bdw - Inches(0.2)) / cols
    chip_h = Inches(0.27)
    for i, name in enumerate(bd_list):
        r, c = divmod(i, cols)
        cx = bdx + Inches(0.1) + c * chip_w
        cy = bdy + Inches(0.4) + r * chip_h
        fill = LIGHT if name.endswith("*") else WHITE
        add_rect(s, cx, cy, chip_w - Inches(0.06),
                 chip_h - Inches(0.04), fill,
                 line=RGBColor(0xCF, 0xD6, 0xE0))
        add_text(s, cx, cy, chip_w - Inches(0.06),
                 chip_h - Inches(0.04), name, size=8.5,
                 color=DARK, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER)
    add_text(s, bdx + Inches(0.1), bdy + bdh - Inches(0.25),
             bdw - Inches(0.2), Inches(0.22),
             "* placeholder BD - created for parity with IPv6 RCC model, no IPv4 subnet yet",
             size=9, color=GREY)

    # App profiles + ESG strip inside VRF-EUR
    appy = bdy + bdh + Inches(0.1)
    appw = bdw / 2 - Inches(0.05)
    panel(s, bdx, appy, appw, Inches(1.25),
          "AppProf-NetCentric-V2", [
              "36 EPGs (EPG-AD-V2, EPG-APP-SVR-V2, ...)",
              "Bound to per-fabric VMMs (APCG-VDS1, APCK-VDS1)",
              "Dynamic VLAN 3501-3967",
          ], accent=CISCO_BLUE)
    panel(s, bdx + appw + Inches(0.1), appy, appw, Inches(1.25),
          "ESG-All-Internal-V2  (Phase 2 in flight)", [
              "Selects all 36 EPGs in AppProf-NetCentric-V2 above",
              "Lives in AppProf-AppCentric-V2 ANP (APIC-direct)",
              "vzAny on VRF-AFR-DEL.Services-V2 keeps it reachability-neutral",
          ], accent=GREEN)

    # VRF-DMZ container
    dx = ex + ew + Inches(0.15)
    dy = ey
    dw = tw - (dx - tx) - Inches(0.2)
    dh = Inches(5.0)
    add_rect(s, dx, dy, dw, dh, LIGHT,
             line=RGBColor(0xCF, 0xD6, 0xE0))
    add_rect(s, dx, dy, dw, Inches(0.45), RED)
    add_text(s, dx + Inches(0.15), dy, dw, Inches(0.45),
             "VRF-DMZ-V2  -  Proxy   |   vzAny: Any_VRF-DMZ-V2",
             size=12, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    dmz_bds = ["BD-D64-PROXY*", "BD-FWEB-PROXY (3)", "BD-RWEB-PROXY*"]
    for i, name in enumerate(dmz_bds):
        bx = dx + Inches(0.15)
        by = dy + Inches(0.55) + i * Inches(0.4)
        fill = LIGHT if name.endswith("*") else WHITE
        add_rect(s, bx, by, dw - Inches(0.3), Inches(0.32),
                 fill, line=RGBColor(0xCF, 0xD6, 0xE0))
        add_text(s, bx, by, dw - Inches(0.3), Inches(0.32),
                 name, size=10, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    panel(s, dx + Inches(0.15), dy + Inches(2.0),
          dw - Inches(0.3), Inches(1.25),
          "AppProf-DMZ-V2",
          ["3 EPGs on per-fabric VMMs", "Dynamic VLAN 3501-3967"],
          accent=CISCO_BLUE)
    panel(s, dx + Inches(0.15), dy + Inches(3.35),
          dw - Inches(0.3), Inches(1.25),
          "ESG-All-DMZ-V2  (Phase 2 in flight)",
          ["Selects all 3 EPGs in AppProf-DMZ-V2 above",
           "Lives in AppProf-AppCentric-V2 ANP (APIC-direct, same as the Internal ESG)"],
          accent=GREEN)

    slides.append(s)


def design_decisions_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Key Design Decisions",
                "Why the target model looks the way it does")
    rows = [
        ["2 VRFs instead of 11",
         "The firewall already enforces real security boundaries - extra VRFs added ops complexity without security value. ESGs + contracts replace VRF-based segmentation."],
        ["VRF-AFR-DEL.Services-V2 (internal)",
         "Consolidates EUR-E (101) + EUR-AIS (132) + EUR-AIM (15) + EUR-AIV (12) + EUR-AIZ (11) + EUR-AIG (1) + EUR-AIP (4) + EUR-GSN-Test (1) ≈ 276 legacy EPGs."],
        ["VRF-DMZ-V2",
         "Keeps EUR-AOV-UC-DMZ and EUR-ARMY-ENT-SVR-DMZ routing-isolated - DMZ must never share a routing table with internal."],
        ["VRF-RCC (IPv6)",
         "Unchanged. Managed separately in ndo-terraform-nac/10.52.4.96/. No IPv6 refactor in this redesign."],
        ["Descriptive naming",
         "BD-DNS-MGMT / EPG-DNS-MGMT replace numeric BD-V0005 / EPG-V0005. Same conventions as the IPv6 RCC model."],
        ["vzAny permit-all (initial)",
         "All EPGs can talk inside a VRF on day 1. ESGs classify endpoints today so contracts can be tightened tomorrow without reclassification."],
        ["L3Out consolidation",
         "Production shrinks 13 L3Outs to ~4 (1 internal + 1 DMZ per site). Lab does not deploy L3Outs."],
        ["L2 'DMZ'-aliased BDs stay in VRF-AFR-DEL.Services-V2",
         "Gateway lives on an external firewall - VRF choice doesn't affect routing. ESG-based grouping (Phase 3) enforces the DMZ/internal boundary instead."],
    ]
    add_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.8),
              ["Decision", "Rationale"], rows,
              col_widths=[1.0, 3.3], body_size=11)
    slides.append(s)


def no_ip_change_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Brownfield-Safe: No IP Addresses Change",
                "Every existing subnet is preserved under the new naming structure")

    # Before / after mapping box
    left_x = Inches(0.5)
    left_w = Inches(6.0)
    box_y = Inches(1.2)
    box_h = Inches(5.6)
    add_rect(s, left_x, box_y, left_w, box_h, LIGHT,
             line=RGBColor(0xDD, 0xE2, 0xEA))
    add_rect(s, left_x, box_y, left_w, Inches(0.45), CISCO_BLUE)
    add_text(s, left_x + Inches(0.15), box_y, left_w, Inches(0.45),
             "Example - BD-APP-SVR-V2 (VRF-AFR-DEL.Services-V2) absorbs 32 subnets",
             size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    mapping = [
        ("10.51.236.65/26", "was BD-V0572 / APCE_EUR_APPS"),
        ("10.51.172.1/24", "was BD-V0372 / APCG_EUR_APPS"),
        ("10.51.108.1/24", "was BD-V0272 / APCK_EUR_APPS"),
        ("10.52.37.129/26", "was BD-V0455 / APCK_RDS"),
        ("…", "28 more subnets absorbed under BD-APP-SVR"),
    ]
    for i, (net, origin) in enumerate(mapping):
        ry = box_y + Inches(0.55) + i * Inches(0.55)
        add_rect(s, left_x + Inches(0.2), ry,
                 left_w - Inches(0.4), Inches(0.48),
                 WHITE, line=RGBColor(0xDD, 0xE2, 0xEA))
        add_text(s, left_x + Inches(0.35), ry,
                 Inches(2.2), Inches(0.48), net,
                 size=13, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left_x + Inches(2.6), ry,
                 left_w - Inches(2.8), Inches(0.48),
                 origin, size=12, color=GREY,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Right column: mechanism + benefits
    right_x = Inches(6.8)
    right_w = Inches(6.1)
    add_text(s, right_x, Inches(1.2), right_w, Inches(0.45),
             "What actually moves", size=16, bold=True, color=NAVY)
    add_bullets(
        s, right_x, Inches(1.7), right_w, Inches(2.5),
        [
            "Gateway MAC/IP stays on ACI - assigned to a renamed BD.",
            "EPG identifiers change, but endpoint IPs don't.",
            "BDs may carry multiple anycast subnets (one BD, many old BDs).",
            "VMM rebinding is transparent - VDS port groups are refreshed by ACI.",
        ],
        size=13,
    )
    add_text(s, right_x, Inches(4.2), right_w, Inches(0.45),
             "Why this de-risks the migration", size=16, bold=True,
             color=NAVY)
    add_bullets(
        s, right_x, Inches(4.7), right_w, Inches(2.3),
        [
            "No DNS, firewall, or NAT changes required.",
            "Hosts/VMs/switches don't get new addresses.",
            "Per-subnet change windows - not a big-bang migration.",
            "Rollback is a VRF reassignment, not a readdressing project.",
        ],
        size=13,
    )
    slides.append(s)


def bd_consolidation_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "BD Consolidation - 215 Legacy → 39 Functional",
                "Every legacy BD is accounted for (mapped, placeholder, or decommissioned)")
    rows = [
        ["BD-APP-SVR", "32 subnets",
         "39 legacy BDs", "Application servers, VDI, ITSM apps"],
        ["BD-CFG-MGMT", "19 subnets",
         "40 legacy BDs", "Server mgmt (OOB, inband, DRAC, IACS)"],
        ["BD-WEB-SVR", "11 subnets",
         "18 legacy BDs", "Web frontends, ITSM web tiers"],
        ["BD-VVOIP-MGMT", "9 subnets",
         "14 legacy BDs", "UC public/restricted, gateways, messaging"],
        ["BD-DB-SVR", "8 subnets",
         "11 legacy BDs", "SQL / database backends"],
        ["BD-MECM", "6 subnets",
         "6 legacy BDs", "SCCM / MECM patching infrastructure"],
        ["BD-BACKUP-SVR", "4 subnets", "— ", "Backup infrastructure"],
        ["BD-FILE-SVR / BD-SYSLOG / BD-LB",
         "3 subnets each", "— ", "File services, logging, load balancers"],
    ]
    add_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(3.2),
              ["New BD", "Subnets", "Legacy BDs absorbed", "Primary function"],
              rows, col_widths=[1.3, 0.8, 1.3, 3.0])

    # Footer split cards
    cards = [
        ("171", "legacy BDs mapped with subnets\n(110 subnets across 22 new BDs)",
         CISCO_BLUE),
        ("14", "L2-only (function known, gateway on firewall)", GREEN),
        ("17", "Placeholder BDs (IPv6-only categories)", ACCENT),
        ("30", "Decommission candidates (20 dead + 4 deprecated + 6 test)",
         RED),
    ]
    cw = Inches(2.95)
    y = Inches(4.7)
    for i, (n, lbl, c) in enumerate(cards):
        x = Inches(0.5) + i * (cw + Inches(0.15))
        add_rect(s, x, y, cw, Inches(2.0), LIGHT,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, x, y, cw, Inches(0.38), c)
        add_text(s, x + Inches(0.15), y, cw - Inches(0.2),
                 Inches(0.38), f"{n} legacy BDs", size=12,
                 bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), y + Inches(0.55),
                 cw - Inches(0.3), Inches(1.4), lbl,
                 size=12, color=DARK)
    slides.append(s)


def naming_conventions_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Naming & Object Model",
                "V2 redesign: -V2 suffix coexists with legacy AFRICOM in the same tenant EUR")
    rows = [
        ["VRF (internal)", "VRF-<scope>-V2", "VRF-AFR-DEL.Services-V2"],
        ["VRF (DMZ)", "VRF-<scope>-V2", "VRF-DMZ-V2"],
        ["VRF (IPv6, legacy)", "VRF-RCC (unchanged)", "VRF-RCC"],
        ["Contract", "Any_<VRF>-V2", "Any_VRF-AFR-DEL.Services-V2, Any_VRF-DMZ-V2"],
        ["Filter", "Any (cross-ref to AFRICOM/VRF_Template/Any)", "Any"],
        ["Bridge Domain", "BD-<function>-V2", "BD-DNS-MGMT-V2, BD-DB-SVR-V2"],
        ["EPG", "EPG-<function>-V2", "EPG-DNS-MGMT-V2, EPG-DB-SVR-V2"],
        ["App Profile (internal EPGs, NDO)", "AppProf-NetCentric-V2",
         "AppProf-NetCentric-V2"],
        ["App Profile (DMZ EPGs, NDO)", "AppProf-DMZ-V2", "AppProf-DMZ-V2"],
        ["App Profile (ESGs, APIC-direct)", "AppProf-AppCentric-V2",
         "AppProf-AppCentric-V2"],
        ["ESG (internal, Phase 2)", "ESG-All-<scope>-V2",
         "ESG-All-Internal-V2"],
        ["ESG (DMZ, Phase 2)", "ESG-All-<scope>-V2", "ESG-All-DMZ-V2"],
        ["ESG (Phase 3, per-zone)", "ESG-<zone>-V2", "ESG-AIM-V2, ESG-DMZ-Web-V2"],
        ["Tenant", "<unchanged>", "EUR"],
        ["Schema", "AFRICOM-V2", "AFRICOM-V2"],
        ["Template", "Tenant_EUR_V2", "Tenant_EUR_V2"],
    ]
    add_table(s, Inches(0.5), Inches(1.15), Inches(8.5), Inches(5.4),
              ["Object", "Pattern", "Example"], rows,
              col_widths=[1.2, 1.6, 1.8])

    panel(
        s, Inches(9.2), Inches(1.15), Inches(3.7), Inches(5.4),
        "Why the -V2 suffix",
        [
            "Legacy AFRICOM schema (production today) and new AFRICOM-V2 schema both deploy into tenant EUR.",
            "ACI requires unique object names per tenant - two NDO templates cannot both own uni/tn-EUR/BD-DB-SVR.",
            "-V2 is GENERATIONAL, not address-family - the same BDs will carry IPv4 + IPv6 (dual-stack) once IPv6 RCC is folded in.",
            "After cutover (AFRICOM retired), the suffix can stay (cosmetic) or be dropped per-object in maintenance windows.",
            "Filter 'Any' is intentionally un-suffixed - cross-referenced from AFRICOM/VRF_Template/Any, not redefined.",
        ],
        accent=ACCENT,
    )
    slides.append(s)


def phased_security_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Phased Security Model",
                "Start permissive, add classification, then tighten - without breaking anyone")
    # Roadmap
    phases = [
        ("Phase 1",
         "vzAny permit-all\non VRF-EUR + VRF-DMZ",
         "Open - any EPG in a VRF can talk",
         "Current", GREEN),
        ("Phase 2",
         "Single ESG per VRF\ngrouping all EPGs",
         "Classification ready, no policy change yet",
         "Current", GREEN),
        ("Phase 3",
         "Split ESGs by zone\n(ESG-AIM, ESG-AIS, ESG-DMZ-Apps, ...)",
         "Segmentation by function with inter-ESG contracts",
         "Future", AMBER),
        ("Phase 4",
         "Tighten contracts\nto required flows only",
         "Micro-segmentation",
         "Future", CISCO_BLUE),
    ]
    box_w = Inches(2.95)
    gap = Inches(0.2)
    start_x = Inches(0.5)
    y = Inches(1.4)
    for i, (name, title, desc, status, color) in enumerate(phases):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, y, box_w, Inches(3.4), LIGHT,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, x, y, box_w, Inches(0.55), color)
        add_text(s, x, y, box_w, Inches(0.55), name, size=14,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), y + Inches(0.7),
                 box_w - Inches(0.4), Inches(1.2), title,
                 size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), y + Inches(2.0),
                 box_w - Inches(0.4), Inches(1.0), desc,
                 size=11, color=DARK, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.7), y + Inches(3.05),
                 box_w - Inches(1.4), Inches(0.3), color)
        add_text(s, x + Inches(0.7), y + Inches(3.05),
                 box_w - Inches(1.4), Inches(0.3), status,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # arrow
        if i < len(phases) - 1:
            ax = x + box_w + Inches(0.02)
            ay = y + Inches(1.5)
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, ay,
                gap - Inches(0.04), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    # Footer takeaway
    add_rect(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.7),
             LIGHT, line=RGBColor(0xDD, 0xE2, 0xEA))
    add_text(s, Inches(0.7), Inches(5.2), Inches(12.0), Inches(0.4),
             "Why phased?", size=14, bold=True, color=NAVY)
    add_bullets(
        s, Inches(0.7), Inches(5.55), Inches(12.0), Inches(1.2),
        [
            "Phase 1+2 deliver the new object model with zero behavior change.",
            "Phases 3-4 add enforcement, one zone / one contract at a time.",
            "Every phase is independently reversible; no phase requires an IP change.",
        ],
        size=12,
    )
    slides.append(s)


def lab_phases_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Lab Migration (Greenfield)",
                "Built from scratch to validate the target design")
    rows = [
        ["1", "Base build",
         "2 VRFs, 39 BDs with consolidated IPv4 subnets, 39 EPGs in AppProf-NetCentric-V2 + AppProf-DMZ-V2, VMM domain, vzAny+permit-all  -  NDO-managed",
         "✓ Complete"],
        ["2", "Lift-and-shift ESGs",
         "ESG-All-Internal-V2 + ESG-All-DMZ-V2 in AppProf-AppCentric-V2 (APIC-direct via nac-aci@0.7.0); EPG-only selectors; vzAny keeps reachability identical",
         "In flight"],
        ["3", "Per-zone ESGs",
         "Split each Phase-2 ESG by adding tag_selectors (vCenter custom-attr) and trimming epg_selectors  -  e.g. ESG-AIM-V2, ESG-DMZ-Web-V2",
         "Future"],
        ["4", "Contract tightening",
         "Replace VRF-level vzAny with explicit ESG-to-ESG contracts on the only flows that need to exist", "Future"],
    ]
    add_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(3.0),
              ["Phase", "Name", "What", "Status"], rows,
              col_widths=[0.4, 1.2, 3.6, 0.8])

    panel(
        s, Inches(0.5), Inches(4.5), Inches(6.0), Inches(2.5),
        "How the lab runs",
        [
            "make init / plan / apply from aci-redesign/apic-vmware/.",
            "nac-aci module (v0.7.0) consumes YAML in data/nac-aci-shared/ + data/nac-aci-site1/.",
            "Local state file; first apply creates ~200+ ACI objects.",
            "Validation: APIC GUI + vCenter port group presence.",
        ],
        accent=CISCO_BLUE,
    )
    panel(
        s, Inches(6.8), Inches(4.5), Inches(6.1), Inches(2.5),
        "Lab tooling",
        [
            "apic-vmware/main.tf - nac-aci module wiring",
            "data/nac-aci-site1/access-policies.nac.yaml - VLAN pool + AAEP (per-fabric)",
            "data/nac-aci-site1-rendered/vmm-domain.nac.yaml - rendered from TF_VAR_vcenter_* (gitignored)",
            "data/nac-aci-shared/tenant-epg-nac.nac.yaml - 2 VRFs, 39 BDs/EPGs, ESGs (cross-fabric)",
        ],
        accent=ACCENT,
    )
    slides.append(s)


def production_phases_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Production Migration (Brownfield)",
                "Coexistence model - live traffic keeps flowing through every phase")
    rows = [
        ["1", "Create VRF-AFR-DEL.Services-V2 and VRF-DMZ-V2 alongside legacy 11 VRFs",
         "None - new VRFs are empty"],
        ["2", "Create 39 new BD-*-V2 with descriptive names in VRF-AFR-DEL.Services-V2/VRF-DMZ-V2",
         "None - new BDs have no endpoints yet"],
        ["3", "Create ESG-*-V2 and vzAny contracts (Any_VRF-*-V2) on VRF-*-V2",
         "None - policy ready, no enforcement yet"],
        ["4", "Migrate EPGs one-at-a-time from old VRFs to VRF-*-V2",
         "Brief traffic loss per subnet - schedule maintenance windows"],
        ["5", "Consolidate L3Outs (13 → ~4)",
         "Routing re-convergence - coordinate with FW/WAN teams"],
        ["6", "Rename EPGs/BDs from numeric to descriptive (-V2)",
         "Cosmetic only - can be done during or after"],
        ["7", "Decommission legacy AFRICOM schema (VRFs, contracts, L3Outs) once empty",
         "None (once all EPGs migrated and validated)"],
    ]
    add_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(4.5),
              ["Phase", "Action", "Key consideration"], rows,
              col_widths=[0.35, 2.4, 2.6], body_size=12)

    add_rect(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.15),
             LIGHT, line=RGBColor(0xDD, 0xE2, 0xEA))
    add_rect(s, Inches(0.5), Inches(5.8), Inches(0.15), Inches(1.15), RED)
    add_text(s, Inches(0.8), Inches(5.85), Inches(12.0), Inches(0.35),
             "Phase 4 is the critical step",
             size=14, bold=True, color=RED)
    add_text(s, Inches(0.8), Inches(6.2), Inches(12.0), Inches(0.75),
             "Moving a BD from one VRF to another causes endpoint re-learning (~seconds of traffic "
             "loss per subnet). Each subnet gets its own change window. Plan with app teams.",
             size=12, color=DARK)
    slides.append(s)


def risks_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Risks & Mitigations",
                "Where traffic loss can occur and how we contain it")
    rows = [
        ["BD VRF reassignment\n(Phase 4)",
         "Endpoint re-learning on the BD - brief L2/L3 blip per subnet.",
         "Per-subnet change windows · notify app teams · rollback = reverse VRF reassignment",
         "HIGH"],
        ["L3Out consolidation\n(Phase 5)",
         "Routing re-convergence at the fabric boundary.",
         "Coordinate with firewall/WAN teams · phased cutover · route monitoring during change",
         "HIGH"],
        ["VMM VDS rebinding",
         "Port group updates in vCenter as EPGs move between VMM-bound BDs.",
         "Use read-write VMM domain · validate in lab first · schedule outside peak hours",
         "MED"],
        ["Contract tightening\n(Phase 3-4)",
         "Flows that weren't designed for can break.",
         "Log-first mode · use ESGs to group before tightening · staged contract rollout",
         "MED"],
        ["Naming / alias drift",
         "Legacy scripts referencing numeric EPG names break.",
         "Keep aliases during Phase 6 · rename with deprecation notice · update automation",
         "LOW"],
        ["Decommission timing",
         "Removing old VRFs/L3Outs before all EPGs migrate.",
         "Gate decom on validation checklist per subnet · automated orphan object report",
         "LOW"],
    ]
    add_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.6),
              ["Risk", "Impact", "Mitigation", "Level"], rows,
              col_widths=[1.0, 1.6, 2.6, 0.4], body_size=11)
    slides.append(s)


def vmm_vlan_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "VMM & VLAN Strategy",
                "Dynamic VLAN assignment under the new VMM domain")

    # VLAN pool bar
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4),
             "VLAN allocation plan", size=16, bold=True, color=NAVY)
    bar_y = Inches(1.65)
    bar_h = Inches(0.7)
    total_range = 4095
    segments = [
        (3001, 3500, "3001-3500", "IPv6 RCC EPGs (static)", CISCO_BLUE),
        (3501, 3967, "3501-3967", "VMM domain (IPv4 EPGs, dynamic)", GREEN),
        (3968, 4095, "3968-4095", "ACI reserved (do not use)", RED),
    ]
    # Use only the region shown (3000-4095)
    start_vlan = 3000
    end_vlan = 4095
    span = end_vlan - start_vlan
    bar_x = Inches(0.5)
    bar_w = Inches(12.3)
    add_rect(s, bar_x, bar_y, bar_w, bar_h, LIGHT,
             line=RGBColor(0xDD, 0xE2, 0xEA))
    for lo, hi, label, desc, color in segments:
        sx = bar_x + bar_w * (lo - start_vlan) / span
        sw = bar_w * (hi - lo + 1) / span
        add_rect(s, sx, bar_y, sw, bar_h, color)
        add_text(s, sx, bar_y, sw, bar_h, label, size=12,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    # Labels beneath
    for lo, hi, label, desc, color in segments:
        sx = bar_x + bar_w * (lo - start_vlan) / span
        sw = bar_w * (hi - lo + 1) / span
        add_text(s, sx, bar_y + bar_h + Inches(0.1), sw,
                 Inches(0.4), desc, size=10, color=color,
                 align=PP_ALIGN.CENTER)

    # Detail tables
    rows = [
        ["VMM domain", "Per-fabric: APCG-VDS1 (Site1), APCK-VDS1 (Site2)"],
        ["vCenter controller", "vcenter01 with credential policy (shared vCenter)"],
        ["Virtual Distributed Switch",
         "Adopted from existing per-fabric VDS (dvs_version='unmanaged')"],
        ["AAEP", "vmm-aaep (lab) / fi-aaep (prod, with phys-fi-domain)"],
        ["Port channel policy", "mac-pin"],
        ["Uplinks", "uplink1, uplink2"],
        ["Leaf interface profile",
         "leaf-101-102-intprof (ports 1-48, nodes 101-102)"],
    ]
    add_table(s, Inches(0.5), Inches(3.3), Inches(7.8), Inches(3.6),
              ["Component", "Configuration"], rows,
              col_widths=[1.0, 2.5])

    panel(
        s, Inches(8.6), Inches(3.3), Inches(4.3), Inches(3.6),
        "Production config touch-points (Design A: UCS-FI direct attach)",
        [
            "data/nac-aci-{site1,site2}-prod/access-policies.nac.yaml - fi-static-vlan-pool (213 VLANs), fi-aaep, PC_FI_A/PC_FI_B policy groups, leaf 101/102 (both sites) split between VMM ports (8-48) and FI uplinks (eth1/6, eth1/7).",
            "data/nac-aci-shared/modules.nac.yaml - aci_mcp sub-module disabled so MCP InstP is managed inline in apic-vmware-prod/main.tf.",
            "data/nac-ndo/schema-africom-v2.nac.yaml - 39 EPG site entries bound to APCG-VDS1 (Site1) and APCK-VDS1 (Site2).",
            "vCenter creds flow in via TF_VAR_vcenter_* env vars (no CHANGE_ME placeholders in tracked YAML).",
        ],
        accent=ACCENT,
    )
    slides.append(s)


def deployed_objects_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "What Gets Deployed",
                "Concrete ACI objects created by terraform apply")

    # Left: access / fabric
    add_text(s, Inches(0.5), Inches(1.15), Inches(6.1), Inches(0.4),
             "Access & Fabric Policies", size=16, bold=True, color=NAVY)
    rows_af = [
        ["VLAN pools",
         "vmm-vlan-pool (dynamic, 3501-3967) + fi-static-vlan-pool (static, 213 VLANs, prod)"],
        ["CDP / LLDP", "cdp-enabled, lldp-enabled (admin_rx_state, admin_tx_state)"],
        ["Port channel", "mac-pin (created from port_channel_policies)"],
        ["Link level", "10G"],
        ["VMware VMM domain",
         "Per-fabric: APCG-VDS1 (Site1), APCK-VDS1 (Site2) - adopt existing VDS"],
        ["vCenter controller", "vcenter01 with credential policy (shared vCenter)"],
        ["VDS uplinks", "uplink1, uplink2"],
        ["AAEP", "vmm-aaep - linked to VMM domain"],
        ["VPC interface policy group", "vpc-vmm-hosts"],
        ["Leaf interface profile",
         "leaf-101-102-intprof (both sites)"],
        ["Leaf switch profile",
         "leaf-101-102-prof (nodes 101-102, both sites)"],
    ]
    add_table(s, Inches(0.5), Inches(1.6), Inches(6.1), Inches(5.3),
              ["Object", "Detail"], rows_af,
              col_widths=[1.0, 2.0], body_size=10)

    # Right: tenant
    add_text(s, Inches(6.9), Inches(1.15), Inches(6.0), Inches(0.4),
             "Tenant EUR", size=16, bold=True, color=NAVY)
    rows_tenant = [
        ["Filter", "Any (cross-ref AFRICOM/VRF_Template/Any, not redefined)"],
        ["Contract", "Any_VRF-AFR-DEL.Services-V2 (scope: context)"],
        ["Contract", "Any_VRF-DMZ-V2 (scope: context)"],
        ["VRF-AFR-DEL.Services-V2", "vzAny provider + consumer of Any_VRF-AFR-DEL.Services-V2"],
        ["VRF-DMZ-V2", "vzAny provider + consumer of Any_VRF-DMZ-V2"],
        ["Bridge Domains (internal)",
         "36 BDs (BD-*-V2) - multi-subnet from legacy consolidation"],
        ["Bridge Domains (DMZ)",
         "3 BDs - BD-D64-PROXY-V2, BD-FWEB-PROXY-V2, BD-RWEB-PROXY-V2"],
        ["App Profile - NetCentric",
         "AppProf-NetCentric-V2 - 36 EPG-*-V2 on per-fabric VMMs (dynamic VLAN)"],
        ["App Profile - DMZ",
         "AppProf-DMZ-V2 - 3 EPG-*-V2 on per-fabric VMMs (dynamic VLAN)"],
        ["App Profile - AppCentric (APIC-direct)",
         "AppProf-AppCentric-V2: ESG-All-Internal-V2 (selects all 36 NetCentric EPGs), ESG-All-DMZ-V2 (selects all 3 DMZ EPGs). Loaded from data/nac-aci-shared/tenant-eur-esgs.nac.yaml -- nac-ndo/mso provider does not model endpoint_security_groups, nac-aci@0.7.0 does."],
    ]
    add_table(s, Inches(6.9), Inches(1.6), Inches(6.0), Inches(5.3),
              ["Object", "Detail"], rows_tenant,
              col_widths=[1.2, 2.3], body_size=10)
    slides.append(s)


def decommission_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Decommission Plan",
                "30 legacy BDs removed during or after migration")
    # Three category columns
    cats = [
        ("Dead / empty (20)", RED,
         ["BD-2250", "BD-V0005", "BD-V0006", "BD-V0958", "BD-V0970",
          "BD-V1116", "BD-V1117", "BD-V1120",
          "BD-V1140 … BD-V1149", "BD-V1571", "BD-V2150"]),
        ("Deprecated (4)", AMBER,
         ["BD-V0009 (Native VLAN)",
          "BD-V0021 (ATM - 'may need to remove')",
          "BD-V0529 ('Remove' in alias)",
          "BD-GSN-Test"]),
        ("Temporary test (6)", GREY,
         ["BD-V0020", "BD-V2001", "BD-V2002",
          "BD-V2003", "BD-V2004", "BD-V2005",
          "(all TMP_SATTest)"]),
    ]
    col_w = Inches(4.0)
    gap = Inches(0.15)
    x0 = Inches(0.5)
    y = Inches(1.2)
    for i, (title, color, items) in enumerate(cats):
        x = x0 + i * (col_w + gap)
        add_rect(s, x, y, col_w, Inches(4.6), LIGHT,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, x, y, col_w, Inches(0.5), color)
        add_text(s, x, y, col_w, Inches(0.5), title, size=14,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.2), y + Inches(0.6),
                    col_w - Inches(0.4), Inches(3.9),
                    items, size=12)

    # Guardrails
    panel(
        s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.95),
        "Decommission guardrails",
        [
            "Only remove after every EPG from an old VRF has migrated and been validated.",
            "Run an orphan-object report (no endpoints, no contracts, no L3Out references) before delete.",
            "Keep an auditable record of each removal - terraform state + APIC audit log.",
        ],
        accent=NAVY,
    )
    slides.append(s)


def what_changes_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "What Changes vs What Stays the Same",
                "Quick sanity check before we start")
    rows = [
        ["IP addresses", "—",
         "All subnets preserved, no readdressing"],
        ["VRF count", "11 → 2", "—"],
        ["BD / EPG count", "215 / 266 → 39 / 39", "—"],
        ["BD / EPG names",
         "Numeric → descriptive", "—"],
        ["Contracts",
         "Per-EPG → vzAny + ESG", "Contract model (still ACI contracts)"],
        ["VMM domain", "VMM1 → APCG-VDS1 / APCK-VDS1",
         "Dynamic VLAN assignment (per-fabric VDS adoption)"],
        ["L3Outs", "13 → ~4 (production)", "External routing concept"],
        ["Firewall", "—",
         "Still enforces DMZ boundaries"],
        ["IPv6 (VRF-RCC)", "—",
         "Unchanged, managed separately"],
    ]
    add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.0),
              ["Aspect", "Changes", "Stays the same"], rows,
              col_widths=[1.0, 1.5, 2.5], body_size=12)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
             "The redesign rewrites the object model, not the network.",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    slides.append(s)


# --- Automation / Deployment slides ---------------------------------------

def deployment_architecture_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Deployment Architecture",
                "How YAML intent becomes ACI configuration")

    # 3 columns: Inputs -> Terraform -> APIC
    col_w = Inches(4.05)
    col_h = Inches(4.3)
    y = Inches(1.2)
    xs = [Inches(0.5), Inches(4.65), Inches(8.8)]
    titles = ["Inputs (source-controlled YAML)",
              "Terraform (orchestration)",
              "APIC (enforcement)"]
    colors = [CISCO_BLUE, ACCENT, GREEN]
    bodies = [
        [
            "data/nac-aci-shared/ - cross-fabric tenant (VRFs, BDs, EPGs, ESGs, contracts)",
            "data/nac-aci-site1/ - per-fabric access policies (Site1)",
            "data/nac-aci-site1-rendered/ - vmm-domain YAML (gitignored)",
            "templates/vmm-domain.nac.yaml.tftpl - template for the above",
            "terraform.tfvars - non-sensitive only (apic_url, username)",
        ],
        [
            "module \"aci\" - netascode/nac-aci@0.7.0 (wrapper)",
            "    reads all YAML files from the two yaml_directories",
            "    aci_mcp sub-module DISABLED via modules.nac.yaml",
            "module \"aci_mcp\" - standalone, owns MCP Instance Policy",
            "    receives sensitive var.mcp_key from env/CI/Vault",
            "provider \"aci\" - CiscoDevNet/aci ≥ 2.0.0",
        ],
        [
            "~200+ objects: VRFs, BDs, EPGs, ESGs, contracts, filters",
            "AppProf-NetCentric-V2 (NDO, 36 EPGs) / AppProf-DMZ-V2 (NDO, 3 EPGs) / AppProf-AppCentric-V2 (APIC-direct, 2 ESGs)",
            "Per-fabric VMM domains (APCG-VDS1, APCK-VDS1) + dynamic VLAN pool 3501-3967",
            "MCP Instance Policy (mcpInstP-default) with compliant key",
            "Access policies: AAEP, domain profile, interface selectors",
        ],
    ]
    for x, title, color, body in zip(xs, titles, colors, bodies):
        add_rect(s, x, y, col_w, col_h, LIGHT,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, x, y, col_w, Inches(0.5), color)
        add_text(s, x, y, col_w, Inches(0.5), title, size=13,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.15), y + Inches(0.6),
                    col_w - Inches(0.3), col_h - Inches(0.7),
                    body, size=10, line_spacing=1.2)

    # Arrows between columns
    for i in range(2):
        ax = xs[i] + col_w + Inches(0.02)
        ay = y + Inches(1.95)
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, ax, ay,
            Inches(0.55), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = NAVY
        arrow.line.fill.background()

    # Pre-step callout: render-vmm-yaml.sh
    add_rect(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(1.25),
             LIGHT, line=RGBColor(0xDD, 0xE2, 0xEA))
    add_rect(s, Inches(0.5), Inches(5.75), Inches(0.15), Inches(1.25), ACCENT)
    add_text(s, Inches(0.8), Inches(5.8), Inches(12.0), Inches(0.35),
             "Pre-step: scripts/render-vmm-yaml.sh  (orchestrated by Makefile / CI)",
             size=13, bold=True, color=NAVY)
    add_text(s, Inches(0.8), Inches(6.15), Inches(12.0), Inches(0.8),
             "Substitutes TF_VAR_vcenter_* env vars into the template and writes the rendered VMM YAML "
             "into the gitignored data/nac-aci-site1-rendered/ directory BEFORE terraform plan runs. "
             "Decoupling the render from Terraform is what lets nac-aci evaluate its for_each/count at plan time.",
             size=11, color=DARK)
    slides.append(s)


def deployment_design_decisions_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Design Decisions (Automation Layer)",
                "Four non-obvious choices; each protects the pipeline from a specific failure mode")
    rows = [
        ["1. aci_mcp sub-module is disabled\nin the nac-aci wrapper",
         "nac-aci unconditionally creates the MCP Instance Policy with a hard-coded default key of 'cisco', which APIC 5.2+/6.x rejects with Error 182.",
         "Disable via modules.nac.yaml (modules.aci_mcp: false). Own the MCP policy directly in main.tf with a sensitive var.mcp_key sourced from env/CI/Vault."],
        ["2. VMM YAML is rendered by a shell\nscript, not a local_file resource",
         "A module with depends_on on an un-applied resource defers every internal for_each/count to apply-time. Result: terraform plan fails with a wall of \"Invalid count argument\" errors.",
         "Render the template in scripts/render-vmm-yaml.sh BEFORE terraform runs. The rendered YAML is a static input at plan time, so nac-aci's graph evaluates cleanly."],
        ["3. Secrets are NOT declared in\nterraform.tfvars",
         "Terraform variable precedence: values in terraform.tfvars beat TF_VAR_* env vars. Leaving an empty slot silently overrides the env var and breaks auth.",
         "Leave apic_password and mcp_key undeclared in tfvars entirely. The env var / CI variable / Vault data source always wins."],
        ["4. vCenter values are NOT\nTerraform variables",
         "Putting the vCenter service-account password into a Terraform variable puts it into state and plan output, even marked sensitive.",
         "Consume them directly in render-vmm-yaml.sh. They never become TF variables, never enter state, never appear in plan output. Rotation = re-export + re-plan."],
    ]
    add_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.6),
              ["Decision", "Problem it solves", "How we implemented it"],
              rows, col_widths=[1.1, 1.8, 2.1], body_size=10,
              header_size=12)

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             "If any of these four are \"fixed\" back to the obvious shape, the pipeline breaks. "
             "Full rationale: aci-redesign/README.md → Design Decisions.",
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    slides.append(s)


def secrets_strategy_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Secrets Strategy: Env → CI → Vault",
                "Same Terraform code; the source of secrets evolves with maturity")

    stages = [
        ("Stage 1: Local Dev", CISCO_BLUE,
         "Shell environment variables",
         [
             "source scripts/set-apic-password.sh",
             "eval \"$(./scripts/generate-mcp-key.sh)\"",
             "export TF_VAR_vcenter_username='...'",
             "export TF_VAR_vcenter_password='...'",
             "",
             "make auth-check → HTTP 200",
             "make plan / make apply",
         ]),
        ("Stage 2: GitLab CI", ACCENT,
         "Masked + Protected CI variables",
         [
             "Settings → CI/CD → Variables:",
             "  APIC_PASSWORD        (masked, protected)",
             "  MCP_KEY              (masked, protected)",
             "  VCENTER_USERNAME     (masked, protected)",
             "  VCENTER_PASSWORD     (masked, protected)",
             "  VCENTER_HOSTNAME_IP  (plain)",
             "  VCENTER_DATACENTER   (plain)",
             "  VCENTER_DVS_VERSION  (plain)",
             ".gitlab-ci.yml maps each to TF_VAR_*.",
         ]),
        ("Stage 3: Vault", GREEN,
         "Centralized KV store (AppRole for runners)",
         [
             "data \"vault_kv_secret_v2\" \"aci_lab\" {",
             "  mount = \"secret\"",
             "  name  = \"aci/lab\"",
             "}",
             "",
             "Wrapper pulls secrets into TF_VAR_* env",
             "at pipeline start. No Terraform code change.",
             "Same lab/prod parity, centralized audit.",
         ]),
    ]
    col_w = Inches(4.05)
    col_h = Inches(5.3)
    y = Inches(1.2)
    xs = [Inches(0.5), Inches(4.65), Inches(8.8)]
    for (title, color, sub, lines), x in zip(stages, xs):
        add_rect(s, x, y, col_w, col_h, LIGHT,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, x, y, col_w, Inches(0.55), color)
        add_text(s, x, y, col_w, Inches(0.55), title,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), y + Inches(0.65),
                 col_w - Inches(0.4), Inches(0.35),
                 sub, size=11, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        # Code-ish body
        tb = s.shapes.add_textbox(x + Inches(0.2),
                                  y + Inches(1.05),
                                  col_w - Inches(0.4),
                                  col_h - Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.2
            r = p.add_run()
            r.text = line if line else " "
            r.font.name = "Consolas"
            r.font.size = Pt(10)
            r.font.color.rgb = DARK

    # Arrows between columns
    for i in range(2):
        ax = xs[i] + col_w + Inches(0.02)
        ay = y + Inches(2.3)
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, ax, ay,
            Inches(0.55), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = NAVY
        arrow.line.fill.background()

    add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
             "The 7 TF_VAR_* names never change. Only the producer changes: shell → CI → Vault.",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    slides.append(s)


def local_workflow_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Local Developer Workflow",
                "What an operator actually types, start to finish")

    # Left: step-by-step commands
    steps = [
        ("1", "Enter the project",
         "cd aci-redesign/apic-vmware"),
        ("2", "Export APIC password\n(hidden prompt, one per shell)",
         "source scripts/set-apic-password.sh"),
        ("3", "Export strong MCP key",
         "eval \"$(./scripts/generate-mcp-key.sh)\""),
        ("4", "Export vCenter values",
         "export TF_VAR_vcenter_hostname_ip='198.18.134.80'\n"
         "export TF_VAR_vcenter_datacenter='Datacenter'\n"
         "export TF_VAR_vcenter_username='administrator@vsphere.local'\n"
         "export TF_VAR_vcenter_password='<single-quoted>'\n"
         "export TF_VAR_vcenter_dvs_version='unmanaged'"),
        ("5", "Prove credentials before touching state",
         "make auth-check     # must print HTTP 200"),
        ("6", "Plan and apply",
         "make plan            # renders VMM YAML + terraform plan\n"
         "make apply           # re-renders + applies plan.tfplan"),
    ]
    y = Inches(1.15)
    x = Inches(0.5)
    w = Inches(7.8)
    step_h = Inches(0.92)
    for num, label, cmd in steps:
        add_rect(s, x, y, Inches(0.55), step_h - Inches(0.08), NAVY)
        add_text(s, x, y, Inches(0.55), step_h - Inches(0.08),
                 num, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x + Inches(0.6), y, w - Inches(0.6),
                 step_h - Inches(0.08), WHITE,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_text(s, x + Inches(0.75), y + Inches(0.03),
                 w - Inches(0.8), Inches(0.3),
                 label, size=11, bold=True, color=NAVY)
        tb = s.shapes.add_textbox(x + Inches(0.75), y + Inches(0.32),
                                  w - Inches(0.8), step_h - Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(cmd.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.1
            r = p.add_run()
            r.text = line
            r.font.name = "Consolas"
            r.font.size = Pt(9)
            r.font.color.rgb = CISCO_BLUE
        y += step_h

    # Right column: context panels
    panel(
        s, Inches(8.6), Inches(1.15), Inches(4.3), Inches(2.4),
        "Why each step matters",
        [
            "Steps 1-4 populate the 7 TF_VAR_* the Makefile checks.",
            "Step 5 isolates credential bugs from Terraform bugs.",
            "Steps 6 renders YAML + runs terraform; make fails fast if any TF_VAR_* is unset.",
        ],
        accent=CISCO_BLUE,
    )
    panel(
        s, Inches(8.6), Inches(3.7), Inches(4.3), Inches(2.0),
        "Env vars don't survive a new terminal",
        [
            "Every new shell repeats steps 2-4.",
            "Use one terminal per session to avoid drift.",
            "set-apic-password.sh must be sourced, not executed.",
        ],
        accent=ACCENT,
    )
    panel(
        s, Inches(8.6), Inches(5.8), Inches(4.3), Inches(1.15),
        "First-time machine setup",
        [
            "make init once per checkout (downloads providers).",
            "Details: aci-redesign/README.md → Getting Started.",
        ],
        accent=GREEN,
    )
    slides.append(s)


def cicd_pipeline_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "CI/CD Pipeline (GitLab)",
                "Four stages; same TF_VAR_* names as local dev, sourced from masked CI variables")

    stages = [
        ("validate", CISCO_BLUE, "on every push",
         [
             "ci-secret-scan.sh (fails pipeline on plaintext secrets)",
             "render-vmm-yaml.sh",
             "terraform fmt -check -recursive",
             "terraform init -backend=false",
             "terraform validate",
         ]),
        ("plan", ACCENT, "on every push",
         [
             "render-vmm-yaml.sh",
             "terraform init + terraform plan -out=plan.tfplan",
             "Artifact: plan.tfplan (30 days)",
             "Needs all 7 TF_VAR_* masked CI variables",
         ]),
        ("deploy", GREEN, "manual gate",
         [
             "terraform apply plan.tfplan",
             "Runs only when a reviewer clicks the play button",
             "Same TF_VAR_* as plan; no surprise credentials",
         ]),
        ("destroy", RED, "manual gate",
         [
             "terraform plan -destroy + terraform apply",
             "Never auto-triggered",
             "Used for tearing down lab experiments cleanly",
         ]),
    ]
    col_w = Inches(3.0)
    gap = Inches(0.1)
    col_h = Inches(4.2)
    y = Inches(1.3)
    start_x = Inches(0.5)
    for i, (name, color, trigger, body) in enumerate(stages):
        x = start_x + i * (col_w + gap)
        add_rect(s, x, y, col_w, col_h, LIGHT,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, x, y, col_w, Inches(0.55), color)
        add_text(s, x, y, col_w, Inches(0.55), name,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + Inches(0.6), col_w, Inches(0.35),
                 trigger, size=11, bold=True, color=color,
                 align=PP_ALIGN.CENTER)
        add_bullets(s, x + Inches(0.15), y + Inches(1.0),
                    col_w - Inches(0.3), col_h - Inches(1.1),
                    body, size=10, line_spacing=1.2)
        if i < 3:
            ax = x + col_w + Inches(0.01)
            ay = y + Inches(1.85)
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, ay,
                gap - Inches(0.02), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = NAVY
            arrow.line.fill.background()

    # Bottom panel: CI var → TF_VAR_ mapping
    add_rect(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(1.25),
             LIGHT, line=RGBColor(0xDD, 0xE2, 0xEA))
    add_rect(s, Inches(0.5), Inches(5.75), Inches(0.15),
             Inches(1.25), NAVY)
    add_text(s, Inches(0.8), Inches(5.78), Inches(12.0), Inches(0.35),
             "GitLab CI variable → TF_VAR_* mapping",
             size=13, bold=True, color=NAVY)
    add_text(s, Inches(0.8), Inches(6.15), Inches(12.0), Inches(0.85),
             "APIC_PASSWORD → TF_VAR_apic_password · MCP_KEY → TF_VAR_mcp_key · "
             "VCENTER_HOSTNAME_IP → TF_VAR_vcenter_hostname_ip · "
             "VCENTER_DATACENTER → TF_VAR_vcenter_datacenter · "
             "VCENTER_USERNAME → TF_VAR_vcenter_username · "
             "VCENTER_PASSWORD → TF_VAR_vcenter_password · "
             "VCENTER_DVS_VERSION → TF_VAR_vcenter_dvs_version",
             size=10, color=DARK)
    slides.append(s)


def ops_troubleshooting_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Operations & Troubleshooting",
                "Small toolset; each command answers one specific question")

    # Left: make auth-check HTTP decision tree
    add_text(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.4),
             "make auth-check  →  interpret the HTTP code",
             size=14, bold=True, color=NAVY)
    codes = [
        ("HTTP 200", "OK", GREEN,
         "Credentials accepted. Proceed with make plan."),
        ("HTTP 401", "Bad credentials", RED,
         "APIC rejected username/password. Re-source set-apic-password.sh; run diagnose-apic-auth.sh to compare env-var password vs freshly-typed."),
        ("HTTP 403", "Account locked", AMBER,
         "Too many failures. Wait 5-10 min, or unlock via another admin (APIC → Admin → AAA → Security → Users)."),
        ("HTTP 000", "No response", RED,
         "APIC unreachable / TLS failure / wrong URL. Test: curl -k https://<apic>/api/aaaLogin.json"),
    ]
    y = Inches(1.6)
    for code, name, color, desc in codes:
        add_rect(s, Inches(0.5), y, Inches(6.0), Inches(1.05),
                 WHITE, line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, Inches(0.5), y, Inches(0.15), Inches(1.05), color)
        add_text(s, Inches(0.75), y + Inches(0.05), Inches(1.6),
                 Inches(0.35), code, size=14, bold=True, color=color)
        add_text(s, Inches(2.35), y + Inches(0.05), Inches(4.0),
                 Inches(0.35), name, size=12, bold=True, color=NAVY)
        add_text(s, Inches(0.75), y + Inches(0.4), Inches(5.6),
                 Inches(0.6), desc, size=10, color=DARK)
        y += Inches(1.15)

    # Right: command reference
    panel(
        s, Inches(6.8), Inches(1.15), Inches(6.1), Inches(2.8),
        "Command reference",
        [
            "make auth-check        credential probe (no terraform)",
            "./scripts/diagnose-apic-auth.sh    env-var vs typed comparison",
            "source scripts/set-apic-password.sh   refresh stale password",
            "make render            re-generate VMM YAML only",
            "make plan / apply      normal change cycle",
            "make destroy           tear everything down",
            "make clean             wipe plan + rendered YAML",
        ],
        accent=CISCO_BLUE,
    )

    panel(
        s, Inches(6.8), Inches(4.05), Inches(6.1), Inches(2.9),
        "Failure modes & fix",
        [
            "terraform plan shows \"Invalid count\" → rendered VMM YAML missing; run make render.",
            "\"Authentication details not provided\" → TF_VAR_apic_password unset; source set-apic-password.sh.",
            "APIC 401 even after re-export → you exported in a DIFFERENT shell than the one running make. Same terminal only.",
            "\"Password is required for MCP Instance Policy\" → TF_VAR_mcp_key missing or < 8 chars; regenerate with generate-mcp-key.sh.",
            "CI secret-scan failure → plaintext password committed to tracked YAML or .tfvars; remove and force-rotate the credential.",
        ],
        accent=RED,
    )
    slides.append(s)


def ndo_split_architecture_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Two-Root Architecture: APIC-direct + NDO",
                "Each control plane drives the layer it actually owns")

    col_w = Inches(6.05)
    col_h = Inches(5.2)
    y = Inches(1.2)
    xs = [Inches(0.5), Inches(6.8)]

    # Left column: apic-vmware/
    add_rect(s, xs[0], y, col_w, col_h, LIGHT,
             line=RGBColor(0xDD, 0xE2, 0xEA))
    add_rect(s, xs[0], y, col_w, Inches(0.55), CISCO_BLUE)
    add_text(s, xs[0], y, col_w, Inches(0.55),
             "apic-vmware/  -  APIC-direct",
             size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, xs[0] + Inches(0.25), y + Inches(0.7),
                col_w - Inches(0.5), col_h - Inches(0.85),
                [
                    "Two providers via aliases: aci (Site1) + aci.site2 (Site2).",
                    "Owns access/fabric policies (leaf profiles, AAEP, VPC).",
                    "Owns MCP Instance Policy per fabric (sensitive key).",
                    "Owns VMware VMM domain (vmmDomP) per fabric.",
                    "manage_tenants = false  -  does NOT create tenant content.",
                    ("Why APIC-direct: leaf profiles and the VMM domain "
                     "object are per-fabric and not modelable in NDO.",
                     1),
                ],
                size=12)

    # Right column: ndo/
    add_rect(s, xs[1], y, col_w, col_h, LIGHT,
             line=RGBColor(0xDD, 0xE2, 0xEA))
    add_rect(s, xs[1], y, col_w, Inches(0.55), NAVY)
    add_text(s, xs[1], y, col_w, Inches(0.55),
             "ndo/  -  NDO-managed",
             size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, xs[1] + Inches(0.25), y + Inches(0.7),
                col_w - Inches(0.5), col_h - Inches(0.85),
                [
                    "Single mso provider (platform=nd, login_domain=local).",
                    "Owns schema AFRICOM-V2 + template Tenant_EUR_V2.",
                    "Owns 2 VRFs (vzAny), 2 contracts, 39 BDs, 2 ANPs, 39 EPGs.",
                    "EPGs bind to APCG-VDS1 (Site1) and APCK-VDS1 (Site2).",
                    "manage_tenants = false  -  EUR pre-exists in NDO.",
                    "deploy_templates = false  -  operator clicks Deploy.",
                    ("Why NDO: only NDO does cross-site policy stitching "
                     "for stretched BDs/EPGs across two APIC fabrics.",
                     1),
                ],
                size=12)

    # Bottom band: cutover order
    band_y = Inches(6.5)
    add_rect(s, Inches(0.5), band_y, Inches(12.3), Inches(0.55), NAVY)
    add_text(s, Inches(0.5), band_y, Inches(12.3), Inches(0.55),
             "Cutover order (clean lab):  apic-vmware/ apply  ->  "
             "ndo/ apply  ->  click Deploy in NDO UI  ->  "
             "scripts/deploy_bindings.py  ->  click Deploy again",
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    slides.append(s)


def ndo_operator_workflow_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "NDO Operator Workflow",
                "Single-template deploy + static port bindings, end-to-end")

    steps = [
        ("1", "Authenticate to NDO",
         "cd aci-redesign/ndo\n"
         "source scripts/set-ndo-password.sh\n"
         "make auth-check     # must return token"),
        ("2", "Apply schema (NDO only, no APIC push)",
         "make plan && make apply\n"
         "# Creates schema AFRICOM-V2 / template Tenant_EUR_V2\n"
         "# 2 VRFs, 2 contracts, 39 BDs, 2 ANPs, 39 EPGs in NDO (all -V2)\n"
         "# deploy_templates=false -> nothing on APICs yet"),
        ("3", "Deploy from NDO UI (manual click)",
         "Schemas -> AFRICOM-V2 -> Tenant_EUR_V2 -> Deploy to sites\n"
         "Review per-site preview -> Confirm\n"
         "After this: tenant content lives on Site1 + Site2"),
        ("4", "Push static port bindings",
         "cd aci-redesign/scripts\n"
         "cp bindings.example.json bindings.json && $EDITOR bindings.json\n"
         "./deploy_bindings.py bindings.json --no-vault --dry-run\n"
         "./deploy_bindings.py bindings.json --no-vault"),
        ("5", "Re-deploy in NDO UI",
         "Static port adds sit pending -> click Deploy again\n"
         "Endpoints can now be plugged in"),
    ]

    y = Inches(1.15)
    x = Inches(0.5)
    w = Inches(12.3)
    step_h = Inches(1.0)
    for num, label, cmd in steps:
        add_rect(s, x, y, Inches(0.55), step_h - Inches(0.08), NAVY)
        add_text(s, x, y, Inches(0.55), step_h - Inches(0.08),
                 num, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x + Inches(0.6), y, w - Inches(0.6),
                 step_h - Inches(0.08), WHITE,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_text(s, x + Inches(0.75), y + Inches(0.04),
                 Inches(4.0), Inches(0.4), label,
                 size=12, bold=True, color=NAVY)
        # command body
        tb = s.shapes.add_textbox(x + Inches(4.9), y + Inches(0.04),
                                  w - Inches(5.05),
                                  step_h - Inches(0.16))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(cmd.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.1
            r = p.add_run()
            r.text = line if line else " "
            r.font.name = "Consolas"
            r.font.size = Pt(10)
            r.font.color.rgb = DARK
        y += step_h

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
             "Day-2 changes: edit data/nac-ndo/, make plan/apply, click Deploy. "
             "Schema refactor: make destroy (NDO-only; APIC untouched), edit, plan, apply, redeploy.",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    slides.append(s)


def next_steps_slide():
    s = prs.slides.add_slide(BLANK)
    page_header(s, "Next Steps",
                "Where this deck hands off")
    columns = [
        ("Lab", GREEN,
         [
             "Phase 1 complete: NDO schema AFRICOM-V2 / single template "
             "Tenant_EUR_V2 deployed via apic-vmware/ + ndo/.",
             "Push static port bindings via scripts/deploy_bindings.py "
             "and re-deploy from NDO UI.",
             "Phase 2 in flight: ESG-All-Internal-V2 + ESG-All-DMZ-V2 in "
             "AppProf-AppCentric-V2, APIC-direct via nac-aci@0.7.0 from "
             "data/nac-aci-shared/tenant-eur-esgs.nac.yaml. EPG-only "
             "selectors; vzAny+permit-all keeps reachability identical.",
             "Phase 3 design ticket: split each Phase-2 ESG into per-zone "
             "ESGs (ESG-AIM-V2, ESG-DMZ-Web-V2, ...) by adding "
             "tag_selectors driven by vCenter custom-attributes.",
         ]),
        ("Production prep", CISCO_BLUE,
         [
             "Review bd_mapping_analysis.txt with app/ops owners.",
             "Confirm L3Out target topology with firewall/WAN teams.",
             "Schedule per-subnet change windows for Phase 4.",
             "Dry-run NDO single-template apply in a non-prod schema.",
         ]),
        ("Tooling / automation", ACCENT,
         [
             "Promote ndo/ root from manual deploy_templates=false "
             "to CI-driven once the runbook is stable.",
             "Add orphan-object report to CI before Phase 7 (decom).",
             "Wire validation checks (EPG counts, BD counts, "
             "VMM bindings) into GitLab CI.",
             "Track placeholder BDs until their IPv4 function is confirmed.",
         ]),
    ]
    col_w = Inches(4.0)
    gap = Inches(0.15)
    x0 = Inches(0.5)
    y = Inches(1.2)
    for i, (title, color, items) in enumerate(columns):
        x = x0 + i * (col_w + gap)
        add_rect(s, x, y, col_w, Inches(5.5), LIGHT,
                 line=RGBColor(0xDD, 0xE2, 0xEA))
        add_rect(s, x, y, col_w, Inches(0.5), color)
        add_text(s, x, y, col_w, Inches(0.5), title,
                 size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.2), y + Inches(0.65),
                    col_w - Inches(0.4), Inches(4.7),
                    items, size=12)

    # Closing line
    add_rect(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
             NAVY)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
             "Same network, cleaner model, safer change path.",
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    slides.append(s)


# --- Build deck ---------------------------------------------------------
title_slide()
agenda_slide()
executive_summary_slide()
current_state_slide()
target_state_slide()
architecture_diagram_slide()
design_decisions_slide()
no_ip_change_slide()
bd_consolidation_slide()
naming_conventions_slide()
phased_security_slide()
lab_phases_slide()
production_phases_slide()
risks_slide()
vmm_vlan_slide()
deployed_objects_slide()
decommission_slide()
what_changes_slide()
deployment_architecture_slide()
deployment_design_decisions_slide()
secrets_strategy_slide()
local_workflow_slide()
cicd_pipeline_slide()
ops_troubleshooting_slide()
ndo_split_architecture_slide()
ndo_operator_workflow_slide()
next_steps_slide()

# Apply footers (skip title)
total = len(slides)
for i, slide in enumerate(slides):
    if i == 0:
        continue
    page_footer(slide, i + 1, total)

out = "/Users/johbarbe/DC/ACI/terraform-esg/aci-redesign/ACI_Redesign_Strategy.pptx"
prs.save(out)
print(f"Saved {out} with {total} slides.")
